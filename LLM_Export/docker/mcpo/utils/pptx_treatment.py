import os
import re
import uuid
import datetime
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt
from pptx.enum.shapes import PP_PLACEHOLDER 
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image as ReportLabImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm
import requests
from utils.file_treatment import _public_url, _generate_unique_folder, _generate_filename, search_image

def _add_table_from_matrix(slide, matrix):
    """
    Create a table on the slide sized to the matrix (rows x cols) and fill it.
    The table is placed over the body placeholder bounds if available,
    else within 1-inch margins.
    Returns the created table shape.
    """
    if not isinstance(matrix, (list, tuple)) or not matrix or not isinstance(matrix[0], (list, tuple)):
        return None

    rows = len(matrix)
    cols = len(matrix[0])

    # determine placement rectangle
    rect = _body_placeholder_bounds(slide)
    if rect:
        left, top, width, height = rect
    else:
        # safe default margins
        left = Inches(1)
        top = Inches(1.2)
        # try to use slide size when available
        try:
            prs = slide.part.presentation
            width = prs.slide_width - Inches(2)
            height = prs.slide_height - Inches(2.2)
        except Exception:
            width = Inches(8)
            height = Inches(4.5)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    # fill cells
    for r in range(rows):
        for c in range(cols):
            try:
                table.cell(r, c).text = "" if matrix[r][c] is None else str(matrix[r][c])
            except Exception:
                pass

    return tbl_shape

def _set_table_from_matrix(shape, data):
    # data = list[list[Any]]; trims to current table size
    tbl = shape.table
    max_r = len(tbl.rows)
    max_c = len(tbl.columns)
    for r, row_vals in enumerate(data):
        if r >= max_r:
            break
        for c, val in enumerate(row_vals):
            if c >= max_c:
                break
            tbl.cell(r, c).text = ""  # clear
            tbl.cell(r, c).text = "" if val is None else str(val)

def _set_text_with_runs(shape, new_content):
    """
    Set the text of a shape while preserving the original run-level formatting.
    """

    if not (hasattr(shape, "text_frame") and shape.text_frame):
        return
    tf = shape.text_frame

    if isinstance(new_content, list):
        lines = [str(item) for item in new_content]
    else:
        lines = [str(new_content or "")]

    original_para_styles = []
    original_para_runs = []  

    for p in tf.paragraphs:
        level = int(getattr(p, "level", 0) or 0)
        alignment = getattr(p, "alignment", None)
        original_para_styles.append({"level": level, "alignment": alignment})
        original_para_runs.append(_snapshot_runs(p))

    tf.clear()

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if (i == 0 and tf.paragraphs) else tf.add_paragraph()

        if original_para_styles:
            style = original_para_styles[i] if i < len(original_para_styles) else original_para_styles[-1]
            p.level = style.get("level", 0)
            if style.get("alignment") is not None:
                p.alignment = style["alignment"]

        runs_spec = (
            original_para_runs[i] if i < len(original_para_runs)
            else (original_para_runs[-1] if original_para_runs else [])
        )

        if not runs_spec:
            r = p.add_run()
            r.text = ""
            continue

        n = len(runs_spec)
        total = len(line)

        if total == 0:
            for spec in runs_spec:
                r = p.add_run()
                r.text = ""
                _apply_font(r, spec["font"])
        else:
            base, rem = divmod(total, n)
            sizes = [base + (1 if k < rem else 0) for k in range(n)]
            pos = 0
            for k, spec in enumerate(runs_spec):
                seg = line[pos:pos + sizes[k]]
                pos += sizes[k]
                r = p.add_run()
                r.text = seg
                _apply_font(r, spec["font"])

def shape_by_id(slide, shape_id):
    sid = int(shape_id)
    for sh in slide.shapes:
        val = getattr(sh, "shape_id", None) or getattr(getattr(sh, "_element", None), "cNvPr", None)
        val = int(getattr(val, "id", val)) if val is not None else None
        if val == sid:
            return sh
    return None

def ensure_slot_textbox(slide, slot):
    slot = (slot or "").lower()

    def _get(ph_name):
        return getattr(PP_PLACEHOLDER, ph_name, None)

    TITLE = _get("TITLE")
    CENTER_TITLE = _get("CENTER_TITLE")
    SUBTITLE = _get("SUBTITLE")
    BODY = _get("BODY")
    CONTENT = _get("CONTENT")
    OBJECT = _get("OBJECT")

    title_types = {t for t in (TITLE, CENTER_TITLE, SUBTITLE) if t is not None}
    body_types  = {t for t in (BODY, CONTENT, OBJECT) if t is not None}

    def find_placeholder(accepted_types):
        for sh in slide.shapes:
            if not getattr(sh, "is_placeholder", False):
                continue
            pf = getattr(sh, "placeholder_format", None)
            if not pf:
                continue
            try:
                if pf.type in accepted_types:
                    return sh
            except Exception:
                pass
        return None

    if slot == "title":
        ph = find_placeholder(title_types)
        if ph:
            return ph

    if slot == "body":
        ph = find_placeholder(body_types)
        if ph:
            return ph

    from pptx.util import Inches
    if slot == "title":
        return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    if slot == "body":
        return slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    return slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))

def _layout_has(layout, want_title=False, want_body=False):
    has_title = has_body = False
    for ph in getattr(layout, "placeholders", []):
        pf = getattr(ph, "placeholder_format", None)
        t = getattr(pf, "type", None) if pf else None
        if t in (getattr(PP_PLACEHOLDER, "TITLE", None),
                 getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
                 getattr(PP_PLACEHOLDER, "SUBTITLE", None)):
            has_title = True
        if t in (getattr(PP_PLACEHOLDER, "BODY", None),
                 getattr(PP_PLACEHOLDER, "CONTENT", None),
                 getattr(PP_PLACEHOLDER, "OBJECT", None)):
            has_body = True
    return (not want_title or has_title) and (not want_body or has_body)

def _pick_layout_for_slots(prs, anchor_slide, needs_title, needs_body):
    if anchor_slide and _layout_has(anchor_slide.slide_layout, needs_title, needs_body):
        return anchor_slide.slide_layout
    for layout in prs.slide_layouts:
        if _layout_has(layout, needs_title, needs_body):
            return layout
    return anchor_slide.slide_layout if anchor_slide else prs.slide_layouts[-1]

def _collect_needs(edit_items):
    needs = {}
    for tgt, _ in edit_items:
        if not isinstance(tgt, str):
            continue
        m = re.match(r"^(n\d+):slot:(title|body)$", tgt.strip(), flags=re.I)
        if m:
            ref, slot = m.group(1), m.group(2).lower()
            needs.setdefault(ref, {"title": False, "body": False})
            needs[ref][slot] = True
    return needs

def _body_placeholder_bounds(slide):
    """Return (left, top, width, height) for the body/content area if possible, else None."""
    try:
        for shp in slide.shapes:
            phf = getattr(shp, "placeholder_format", None)
            if phf is not None:
                # BODY placeholder is the content region on most layouts
                if str(getattr(phf, "type", "")).endswith("BODY"):
                    return shp.left, shp.top, shp.width, shp.height
    except Exception:
        pass
    return None

def _resolve_donor_simple(order, slides_by_id, anchor_id, kind):
    """
    kind: 'insert_after' or 'insert_before'
    Rules:
      insert_after(anchor):
        - if anchor is first -> donor = next slide if exists else anchor
        - else               -> donor = anchor
      insert_before(anchor):
        - if anchor is last  -> donor = previous slide if exists else anchor
        - else               -> donor = anchor
    """
    if not order:
        return None
    if anchor_id not in order:
        # anchor not found
        return slides_by_id.get(order[1]) if len(order) > 1 else slides_by_id.get(order[0])

    pos = order.index(anchor_id)
    last_idx = len(order) - 1

    if kind == "insert_after":
        if pos == 0:
            # after first slide
            return slides_by_id.get(order[pos + 1]) if pos + 1 <= last_idx else slides_by_id.get(anchor_id)
        else:
            return slides_by_id.get(anchor_id)

    # insert_before
    if pos == last_idx:
        # before last slide
        return slides_by_id.get(order[pos - 1]) if pos - 1 >= 0 else slides_by_id.get(anchor_id)
    else:
        return slides_by_id.get(anchor_id)

def _snapshot_runs(p):
    """Return a list of {'text': str, 'font': {...}} for each run in a paragraph."""
    runs = []
    for r in p.runs:
        f = r.font
        font_spec = {
            "name": f.name,
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "color_rgb": getattr(getattr(f.color, "rgb", None), "rgb", None) or getattr(f.color, "rgb", None)
        }
        runs.append({"text": r.text or "", "font": font_spec})
    return runs

def _apply_font(run, font_spec):
    """Apply font specifications to a run."""
    if not font_spec:
        return
    f = run.font
    try:
        if font_spec.get("name") is not None:
            f.name = font_spec["name"]
        if font_spec.get("size") is not None:
            f.size = font_spec["size"]
        if font_spec.get("bold") is not None:
            f.bold = font_spec["bold"]
        if font_spec.get("italic") is not None:
            f.italic = font_spec["italic"]
        if font_spec.get("underline") is not None:
            f.underline = font_spec["underline"]
        rgb = font_spec.get("color_rgb")
        if rgb is not None:
            try:
                f.color.rgb = rgb
            except Exception:
                pass
    except Exception:
        pass

def _get_pptx_namespaces():
    """Returns XML namespaces for PowerPoint"""
    return {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p15': 'http://schemas.microsoft.com/office/powerpoint/2012/main',
        'p14': 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    }

def _add_native_pptx_comment_zip(pptx_path, slide_num, comment_text, author_id, x=100, y=100):
    """
    Add a native PowerPoint comment by directly manipulating the ZIP file.
        Args:
        pptx_path: Path to the PPTX file
        slide_num: Slide number (1-based)
        comment_text: Comment text
        author_id: Author ID
        x: X position in EMU (not pixels!)
        y: Y position in EMU (not pixels!)
    """
    namespaces = _get_pptx_namespaces()
    
    import tempfile
    from pathlib import Path
    import zipfile
    from lxml import etree
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_path)
        
        authors_file = temp_path / 'ppt' / 'commentAuthors.xml'
        if authors_file.exists():
            root = etree.parse(str(authors_file)).getroot()
            found = False
            for author in root.findall('.//p:cmAuthor', namespaces):
                if author.get('name') == 'AI Reviewer':
                    author_id = int(author.get('id'))
                    found = True
                    break
            
            if not found:
                existing_ids = [int(a.get('id')) for a in root.findall('.//p:cmAuthor', namespaces)]
                author_id = max(existing_ids) + 1 if existing_ids else 0
                author = etree.SubElement(root, f'{{{namespaces["p"]}}}cmAuthor')
                author.set('id', str(author_id))
                author.set('name', 'AI Reviewer')
                author.set('initials', 'AI')
                author.set('lastIdx', '1')
                author.set('clrIdx', str(author_id % 8))
        else:
            authors_file.parent.mkdir(parents=True, exist_ok=True)
            root = etree.Element(
                f'{{{namespaces["p"]} }}cmAuthorLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            author = etree.SubElement(root, f'{{{namespaces["p"]} }}cmAuthor')
            author.set('id', str(author_id))
            author.set('name', 'AI Reviewer')
            author.set('initials', 'AI')
            author.set('lastIdx', '1')
            author.set('clrIdx', '0')
            
            rels_file = temp_path / 'ppt' / '_rels' / 'presentation.xml.rels'
            if rels_file.exists():
                rels_root = etree.parse(str(rels_file)).getroot()
            else:
                rels_file.parent.mkdir(parents=True, exist_ok=True)
                rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/commentAuthors')
            rel.set('Target', 'commentAuthors.xml')
            
            with open(rels_file, 'wb') as f:
                f.write(etree.tostring(rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with open(authors_file, 'wb') as f:
            f.write(etree.tostring(root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        comments_dir = temp_path / 'ppt' / 'comments'
        comments_dir.mkdir(parents=True, exist_ok=True)
        comment_file = comments_dir / f'comment{slide_num}.xml'
        
        if comment_file.exists():
            comments_root = etree.parse(str(comment_file)).getroot()
        else:
            comments_root = etree.Element(
                f'{{{namespaces["p"]} }}cmLst',
                nsmap={k: v for k, v in namespaces.items() if k in ['p']}
            )
            
            slide_rels_file = temp_path / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'
            if slide_rels_file.exists():
                slide_rels_root = etree.parse(str(slide_rels_file)).getroot()
            else:
                slide_rels_file.parent.mkdir(parents=True, exist_ok=True)
                slide_rels_root = etree.Element(
                    '{http://schemas.openxmlformats.org/package/2006/relationships}Relationships'
                )
            
            existing_ids = [int(rel.get('Id')[3:]) for rel in slide_rels_root.findall('.//{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')]
            next_rid = max(existing_ids) + 1 if existing_ids else 1
            
            rel = etree.SubElement(slide_rels_root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            rel.set('Id', f'rId{next_rid}')
            rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments')
            rel.set('Target', f'../comments/comment{slide_num}.xml')
            
            with open(slide_rels_file, 'wb') as f:
                f.write(etree.tostring(slide_rels_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        existing_ids = [int(c.get('idx')) for c in comments_root.findall('.//p:cm', namespaces)]
        next_id = max(existing_ids) + 1 if existing_ids else 1
        
        comment = etree.SubElement(comments_root, f'{{{namespaces["p"]} }}cm')
        comment.set('authorId', str(author_id))
        comment.set('dt', datetime.datetime.now().isoformat())
        comment.set('idx', str(next_id))
        
        pos = etree.SubElement(comment, f'{{{namespaces["p"]} }}pos')
        pos.set('x', str(int(x)))
        pos.set('y', str(int(y)))
        
        text_elem = etree.SubElement(comment, f'{{{namespaces["p"]} }}text')
        text_elem.text = comment_text
        
        with open(comment_file, 'wb') as f:
            f.write(etree.tostring(comments_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        content_types_file = temp_path / '[Content_Types].xml'
        if content_types_file.exists():
            ct_root = etree.parse(str(content_types_file)).getroot()
            ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
            
            has_authors = False
            has_comments = False
            
            for override in ct_root.findall('.//ct:Override', ns):
                if override.get('PartName') == '/ppt/commentAuthors.xml':
                    has_authors = True
                if override.get('PartName') == f'/ppt/comments/comment{slide_num}.xml':
                    has_comments = True
            
            if not has_authors:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', '/ppt/commentAuthors.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.commentAuthors+xml')
            
            if not has_comments:
                override = etree.SubElement(ct_root, '{http://schemas.openxmlformats.org/package/2006/content-types}Override')
                override.set('PartName', f'/ppt/comments/comment{slide_num}.xml')
                override.set('ContentType', 'application/vnd.openxmlformats-officedocument.presentationml.comments+xml')
            
            with open(content_types_file, 'wb') as f:
                f.write(etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', pretty_print=True))
        
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in temp_path.rglob('*'):
                if file_path.is_file():
                    arcname = str(file_path.relative_to(temp_path))
                    zf.write(file_path, arcname)
        
        logging.debug(f"Native comment added to slide {slide_num} with idx={next_id}")

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    """
    Calcule une taille de police approximative en fonction de la quantité de texte à rendre.
    Retourne un objet Pt (pptx.util.Pt).
    """
    try:
        total_chars = sum(len(str(line)) for line in (content_list or []))
    except Exception:
        total_chars = 0
    ratio = (total_chars / max_chars) if max_chars and max_chars > 0 else 1
    if ratio <= 1:
        return PptPt(base_size)
    new_size = int(base_size / ratio)
    return PptPt(max(min_size, new_size))

def create_presentation(
    slides_data: list[dict],
    filename: str,
    folder_path: str | None = None,
    title: str | None = None,
    pptx_template_path: str | None = None
) -> dict:
    """
    Crée une présentation PPTX et retourne {'url','path'}.
    - slides_data: liste de dicts {"title": str, "content": list[str|Any], "image_query": str|None, "image_position": str, "image_size": str}
    - filename: nom du fichier à créer (ex: 'slides.pptx')
    - folder_path: dossier de sortie (si None, un dossier unique est généré)
    - title: titre de la présentation (pour la diapo de titre)
    - pptx_template_path: chemin d'un template .pptx optionnel
    """
    if folder_path is None:
        folder_path = _generate_unique_folder()
    if filename:
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)
        fname = filename
    else:
        filepath, fname = _generate_filename(folder_path, "pptx")

    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    # Chargement du template si fourni
    if pptx_template_path and os.path.exists(pptx_template_path):
        try:
            tmp = Presentation(pptx_template_path)
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True
                # Heuristique: on récupère les layouts de la (ou des) premières diapos du template
                try:
                    title_layout = prs.slides[0].slide_layout
                except Exception:
                    title_layout = prs.slide_layouts[0]
                try:
                    content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout
                except Exception:
                    content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

                # Nettoyer toutes les diapos du template sauf la première (pour garder le style)
                try:
                    for i in range(len(prs.slides) - 1, 0, -1):
                        rId = prs.slides._sldIdLst[i].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[i]
                except Exception:
                    pass
        except Exception as e:
            logging.error(f"Error loading PPTX template '{pptx_template_path}': {e}")
            use_template = False
            prs = None

    if not use_template:
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    # Diapo de titre
    if use_template:
        tslide = prs.slides[0]
        if getattr(tslide.shapes, "title", None):
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        title_info = next(
                            (
                                {
                                    'size': PptPt(int(child.attrib.get('sz', 2800)) / 100),
                                    'bold': child.attrib.get('b', '0') == '1'
                                }
                                for child in title_layout.element.iter()
                                if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib
                            ),
                            {'size': PptPt(28), 'bold': True}
                        )
                        r.font.size = title_info['size']
                        r.font.bold = title_info['bold']
                    except Exception:
                        r.font.size = PptPt(28)
                        r.font.bold = True
    else:
        tslide = prs.slides.add_slide(title_layout)
        if getattr(tslide.shapes, "title", None):
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = PptPt(28)
                    r.font.bold = True

    # Constantes de mesures
    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN

    page_margin = 0.5
    gutter = 0.3

    # Ajout des diapositives de contenu
    for i, slide_data in enumerate(slides_data or []):
        if not isinstance(slide_data, dict):
            continue

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]

        slide = prs.slides.add_slide(content_layout)

        # Titre de la diapo
        if getattr(slide.shapes, "title", None):
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        title_info = next(
                            (
                                {
                                    'size': PptPt(int(child.attrib.get('sz', 2800)) / 100),
                                    'bold': child.attrib.get('b', '0') == '1'
                                }
                                for child in content_layout.element.iter()
                                if 'defRPr' in child.tag.split('}')[-1] and 'sz' in child.attrib
                            ),
                            {'size': PptPt(28), 'bold': True}
                        )
                        r.font.size = title_info['size']
                        r.font.bold = title_info['bold']
                    except Exception:
                        r.font.size = PptPt(28)
                        r.font.bold = True

        # Trouver/Créer une zone de contenu
        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph
                        break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph
                            break
                    except Exception:
                        pass
        except Exception as e:
            logging.error(f"Error finding content placeholder: {e}")

        title_bottom_in = 1.0
        if getattr(slide.shapes, "title", None):
            try:
                title_bottom_emu = slide.shapes.title.top + slide.shapes.title.height
                title_bottom_in = max(title_bottom_emu / EMU_PER_IN, 1.0) + 0.2
            except Exception:
                title_bottom_in = 1.2

        if content_shape is None:
            content_shape = slide.shapes.add_textbox(
                Inches(page_margin),
                Inches(title_bottom_in),
                Inches(slide_w_in - 2 * page_margin),
                Inches(slide_h_in - title_bottom_in - page_margin),
            )
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            pass
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        # Zone de texte par défaut
        content_left_in, content_top_in = page_margin, title_bottom_in
        content_width_in = slide_w_in - 2 * page_margin
        content_height_in = slide_h_in - (title_bottom_in + page_margin)

        # Image optionnelle (via query)
        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                try:
                    response = requests.get(image_url, timeout=30)
                    response.raise_for_status()
                    image_data = response.content
                    image_stream = BytesIO(image_data)

                    pos = slide_data.get("image_position", "right")
                    size = slide_data.get("image_size", "medium")
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = title_bottom_in
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = title_bottom_in
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2 * page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = slide_w_in - 2 * page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = title_bottom_in
                        content_left_in = page_margin
                        content_top_in = title_bottom_in
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (title_bottom_in + page_margin)

                    slide.shapes.add_picture(
                        image_stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in)
                    )
                except Exception:
                    pass

        # Positionner/redimensionner la zone de texte
        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        # Estimation de capacité et taille de police
        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        safe_width = max(content_width_in, 0.1)
        safe_height = max(content_height_in, 0.1)
        est_capacity = int(safe_width * approx_chars_per_in * safe_height * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        # Remplir le texte
        try:
            tf = content_shape.text_frame
        except Exception:
            tf = getattr(content_shape, "text_frame", None)

        if tf is None:
            continue

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = str(line) if line is not None else ""
            try:
                run.font.size = font_size
            except Exception:
                pass
            try:
                p.space_after = PptPt(6)
            except Exception:
                pass

    prs.save(filepath)
    return {"url": _public_url(folder_path, fname), "path": filepath}
