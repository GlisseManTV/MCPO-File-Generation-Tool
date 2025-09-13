import re
import os
import ast
import json
import uuid
import emoji
import time
import datetime
import tarfile
import zipfile
import py7zr
import logging
import requests
import threading
import markdown2
import tempfile
from bs4 import BeautifulSoup, NavigableString
from mcp.server.fastmcp import FastMCP
from openpyxl import Workbook
import csv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.parts.image import Image
from pptx.enum.text import MSO_AUTO_SIZE
from io import BytesIO
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm

PERSISTENT_FILES = os.getenv("PERSISTENT_FILES", "false")
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)


BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

PPTX_TEMPLATE_PATH= os.getenv("PPTX_TEMPLATE_PATH",None)
PPTX_TEMPLATE = None
if PPTX_TEMPLATE_PATH and os.path.exists(PPTX_TEMPLATE_PATH):
    PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
    logging.info(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")

def search_image(query):
    api_key = os.getenv("UNSPLASH_ACCESS_KEY")
    if not api_key:
        log.warning("UNSPLASH_ACCESS_KEY is not set. Cannot search for images.")
        return None
    url = "https://api.unsplash.com/search/photos"
    params = {
        "query": query,
        "per_page": 1,
        "orientation": "landscape"
    }
    headers = {"Authorization": f"Client-ID {api_key}"}
    log.debug(f"Searching Unsplash for query: '{query}'")
    try:
        response = requests.get(url, params=params, headers=headers)
        log.debug(f"Unsplash API response status: {response.status_code}")
        response.raise_for_status() 
        data = response.json()
        if data.get("results"):
            image_url = data["results"][0]["urls"]["regular"]
            log.debug(f"Found image URL for '{query}': {image_url}")
            return image_url
        else:
            log.debug(f"No results found on Unsplash for query: '{query}'")
    except requests.exceptions.RequestException as e:
        log.error(f"Network error while searching image for '{query}': {e}")
    except json.JSONDecodeError as e:
        log.error(f"Error decoding JSON from Unsplash for '{query}': {e}")
    except Exception as e:
        log.error(f"Unexpected error searching image for '{query}': {e}")
    return None

def _resolve_log_level(val: str | None) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(
    level=_resolve_log_level(LOG_LEVEL_ENV),
    format=LOG_FORMAT_ENV,
)
log = logging.getLogger("file_export_mcp")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

mcp = FastMCP("file_export")

def dynamic_font_size(content_list, max_chars=400, base_size=28, min_size=12):
    total_chars = sum(len(line) for line in content_list)
    ratio = total_chars / max_chars if max_chars > 0 else 1
    if ratio <= 1:
        return Pt(base_size)
    else:
        new_size = int(base_size / ratio)
        return Pt(max(min_size, new_size))


def _public_url(folder_path: str, filename: str) -> str:
    """Build a stable public URL for a generated file."""
    folder = os.path.basename(folder_path).lstrip("/")
    name = filename.lstrip("/")
    return f"{BASE_URL}/{folder}/{name}"

def _generate_unique_folder() -> str:
    folder_name = f"export_{uuid.uuid4().hex[:10]}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}"
    folder_path = os.path.join(EXPORT_DIR, folder_name)
    os.makedirs(folder_path, exist_ok=True)
    return folder_path

def _generate_filename(folder_path: str, ext: str, filename: str = None) -> tuple[str, str]:
    if not filename:
        filename = f"export_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    return filepath, filename

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(
    name="CustomHeading1",
    parent=styles["Heading1"],
    textColor=colors.HexColor("#0A1F44"),
    fontSize=18,
    spaceAfter=16,
    spaceBefore=12,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading2",
    parent=styles["Heading2"],
    textColor=colors.HexColor("#1C3F77"),
    fontSize=14,
    spaceAfter=12,
    spaceBefore=10,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomHeading3",
    parent=styles["Heading3"],
    textColor=colors.HexColor("#3A6FB0"), 
    fontSize=12,
    spaceAfter=10,
    spaceBefore=8,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomNormal",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomListItem",
    parent=styles["Normal"],
    fontSize=11,
    leading=14,
    alignment=TA_LEFT
))
styles.add(ParagraphStyle(
    name="CustomCode",
    parent=styles["Code"],
    fontSize=10,
    leading=12,
    fontName="Courier",
    backColor=colors.HexColor("#F5F5F5"),
    borderColor=colors.HexColor("#CCCCCC"),
    borderWidth=1,
    leftIndent=10,
    rightIndent=10,
    topPadding=5,
    bottomPadding=5
))

def render_text_with_emojis(text: str) -> str:
    if not text:
        return ""
    try:
        converted = emoji.emojize(text, language="alias")
        return converted
    except Exception as e:
        log.error(f"Error in emoji conversion: {e}")
        return text

def process_list_items(ul_or_ol_element, is_ordered=False):
    items = []
    bullet_type = '1' if is_ordered else 'bullet'
    for li in ul_or_ol_element.find_all('li', recursive=False):
        li_text_parts = []
        for content in li.contents:
            if isinstance(content, NavigableString):
                li_text_parts.append(str(content))
            elif content.name not in ['ul', 'ol']:
                 li_text_parts.append(content.get_text())
        li_text = ''.join(li_text_parts).strip()
        list_item_paragraph = None
        if li_text:
            rendered_text = render_text_with_emojis(li_text)
            list_item_paragraph = Paragraph(rendered_text, styles["CustomListItem"])
        sub_lists = li.find_all(['ul', 'ol'], recursive=False)
        sub_flowables = []
        if list_item_paragraph:
             sub_flowables.append(list_item_paragraph)
        for sub_list in sub_lists:
            is_sub_ordered = sub_list.name == 'ol'
            nested_items = process_list_items(sub_list, is_sub_ordered)
            if nested_items:
                nested_list_flowable = ListFlowable(
                    nested_items,
                    bulletType='1' if is_sub_ordered else 'bullet',
                    leftIndent=10 * mm,
                    bulletIndent=5 * mm,
                    spaceBefore=2,
                    spaceAfter=2
                )
                sub_flowables.append(nested_list_flowable)
        if sub_flowables:
            items.append(ListItem(sub_flowables))
    return items

def render_html_elements(soup):
    log.debug("Starting render_html_elements...")
    story = []
    element_count = 0
    for elem in soup.children:
        element_count += 1
        log.debug(f"Processing element #{element_count}: {type(elem)}, name={getattr(elem, 'name', 'NavigableString')}")
        if isinstance(elem, NavigableString):
            text = str(elem).strip()
            if text:
                log.debug(f"Adding Paragraph from NavigableString: {text[:50]}...")
                story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                story.append(Spacer(1, 6))
        elif hasattr(elem, 'name'):
            tag_name = elem.name
            log.debug(f"Handling tag: <{tag_name}>")
            if tag_name == "h1":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H1: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading1"]))
                story.append(Spacer(1, 10))
            elif tag_name == "h2":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H2: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading2"]))
                story.append(Spacer(1, 8))
            elif tag_name == "h3":
                text = render_text_with_emojis(elem.get_text().strip())
                log.debug(f"Adding H3: {text[:50]}...")
                story.append(Paragraph(text, styles["CustomHeading3"]))
                story.append(Spacer(1, 6))
            elif tag_name == "p":
                imgs = elem.find_all("img")
                if imgs:
                    for img_tag in imgs:
                        src = img_tag.get("src")
                        alt = img_tag.get("alt", "[Image]")
                        try:
                            if src and src.startswith("http"):
                                response = requests.get(src)
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = Image(img_data, width=200, height=150)
                            else:
                                img = Image(src, width=200, height=150)
                            story.append(img)
                            story.append(Spacer(1, 10))
                        except Exception as e:
                            log.error(f"Error loading image {src}: {e}")
                            story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                            story.append(Spacer(1, 6))
                else:
                    text = render_text_with_emojis(elem.get_text().strip())
                    if text:
                        log.debug(f"Adding Paragraph: {text[:50]}...")
                        story.append(Paragraph(text, styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
            elif tag_name in ["ul", "ol"]:
                is_ordered = tag_name == "ol"
                log.debug(f"Processing list (ordered={is_ordered})...")
                items = process_list_items(elem, is_ordered)
                if items:
                    log.debug(f"Adding ListFlowable with {len(items)} items")
                    story.append(ListFlowable(items,
                        bulletType='1' if is_ordered else 'bullet',
                        leftIndent=10 * mm,
                        bulletIndent=5 * mm,
                        spaceBefore=6,
                        spaceAfter=10
                    ))
            elif tag_name == "blockquote":
                text = render_text_with_emojis(elem.get_text().strip())
                if text:
                    log.debug(f"Adding Blockquote: {text[:50]}...")
                    story.append(Paragraph(f"{text}", styles["CustomNormal"]))
                    story.append(Spacer(1, 8))
            elif tag_name in ["code", "pre"]:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Code/Pre block: {text[:50]}...")
                    story.append(Paragraph(text, styles["CustomCode"]))
                    story.append(Spacer(1, 6 if tag_name == "code" else 8))
            elif tag_name == "img":
                src = elem.get("src")
                alt = elem.get("alt", "[Image]")
                log.debug(f"Found <img> tag. src='{src}', alt='{alt}'")
                if src is not None: 
                    try:
                        if src.startswith("image_query:"):

                            query = src.replace("image_query:", "").strip()
                            log.debug(f"Handling image_query: '{query}'")
                            image_url = search_image(query)
                            if image_url:
                                log.debug(f"Downloading image from Unsplash URL: {image_url}")
                                response = requests.get(image_url)
                                log.debug(f"Image download response status: {response.status_code}")
                                response.raise_for_status()
                                img_data = BytesIO(response.content)
                                img = ReportLabImage(img_data, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Unsplash)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                                log.warning(f"No image found for query: {query}")
                                story.append(Paragraph(f"[Image non trouvee pour: {query}]", styles["CustomNormal"]))
                                story.append(Spacer(1, 6))
                        elif src.startswith("http"):
                            log.debug(f"Downloading image from direct URL: {src}")
                            response = requests.get(src)
                            log.debug(f"Image download response status: {response.status_code}")
                            response.raise_for_status()
                            img_data = BytesIO(response.content)
                            img = ReportLabImage(img_data, width=200, height=150)
                            log.debug("Adding ReportLab Image object to story (Direct URL)")
                            story.append(img)
                            story.append(Spacer(1, 10))
                        else:
                            log.debug(f"Loading local image: {src}")
                            if os.path.exists(src):
                                img = ReportLabImage(src, width=200, height=150)
                                log.debug("Adding ReportLab Image object to story (Local)")
                                story.append(img)
                                story.append(Spacer(1, 10))
                            else:
                               log.error(f"Local image file not found: {src}")
                               story.append(Paragraph(f"[Image locale non trouvee: {src}]", styles["CustomNormal"]))
                               story.append(Spacer(1, 6))
                    except requests.exceptions.RequestException as e:
                        log.error(f"Network error loading image {src}: {e}")
                        story.append(Paragraph(f"[Image (erreur reseau): {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                    except Exception as e:
                        log.error(f"Error processing image {src}: {e}", exc_info=True) 
                        story.append(Paragraph(f"[Image: {alt}]", styles["CustomNormal"]))
                        story.append(Spacer(1, 6))
                else:
                    log.warning("Image tag found with no 'src' attribute.")
                    story.append(Paragraph(f"[Image: {alt} (source manquante)]", styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
            elif tag_name == "br":
                log.debug("Adding Spacer for <br>")
                story.append(Spacer(1, 6))
            else:
                text = elem.get_text().strip()
                if text:
                    log.debug(f"Adding Paragraph for unknown tag <{tag_name}>: {text[:50]}...")
                    story.append(Paragraph(render_text_with_emojis(text), styles["CustomNormal"]))
                    story.append(Spacer(1, 6))
    log.debug(f"Finished render_html_elements. Story contains {len(story)} elements.")
    return story

def _cleanup_files(folder_path: str, delay_minutes: int):
    def delete_files():
        time.sleep(delay_minutes * 60)
        try:
            import shutil
            shutil.rmtree(folder_path) 
            log.debug(f"Folder {folder_path} deleted.")
        except Exception as e:
            logging.error(f"Error deleting files : {e}")
    thread = threading.Thread(target=delete_files)
    thread.start()

@mcp.tool()
def create_excel(data: list[list[str]], filename: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    folder_path = _generate_unique_folder()
    filepath, fname = _generate_filename(folder_path, "xlsx", filename)
    wb = Workbook()
    ws = wb.active
    for row in data:
        ws.append(row)
    wb.save(filepath)
    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    return {"url": _public_url(folder_path, fname)}

@mcp.tool()
def create_csv(data: list[list[str]], filename: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    folder_path = _generate_unique_folder()
    filepath, fname = _generate_filename(folder_path, "csv", filename)
    with open(filepath, "w", newline="", encoding="utf-8") as f:
        csv.writer(f).writerows(data)
    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    return {"url": _public_url(folder_path, fname)}

@mcp.tool()
def create_pdf(text: list[str], filename: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    log.debug("Starting create_pdf tool...")
    folder_path = _generate_unique_folder()
    filepath, fname = _generate_filename(folder_path, "pdf", filename)
    md_text = "\n".join(text)
    log.debug(f"Input Markdown text:\n{md_text}")

    def replace_image_query(match):
        query = match.group(1).strip()
        log.debug(f"Found image_query placeholder: '{query}'")
        image_url = search_image(query)

        if image_url:
            result_tag = f'\n\n<img src="{image_url}" alt="Searched image: {query}" />\n\n'
            log.debug(f"Replaced image_query '{query}' with URL: {image_url}")
        else:
            result_tag = f'\n\n<p>[Image not found for: {query}]</p>\n\n'
            log.warning(f"Failed to find image for query: '{query}'")

        log.debug(f"Replacement result: {result_tag}")
        return result_tag

    log.debug("Applying image_query regex replacement...")
    md_text_before_replace = md_text
    md_text = re.sub(r'!\[[^\]]*\]\(\s*image_query:\s*([^)]+)\)', replace_image_query, md_text)
    if md_text != md_text_before_replace:
        log.debug(f"Markdown text after replacement:\n{md_text}")
    else:
        log.debug("No image_query replacements were made.")

    log.debug("Converting Markdown to HTML...")
    html = markdown2.markdown(
        md_text,
        extras=[
            'fenced-code-blocks',
            'tables',
            'break-on-newline',
            'cuddled-lists', 
            'smarty-pants'
        ]
    )
    log.debug(f"Generated HTML:\n{html}") 

    log.debug("Parsing HTML with BeautifulSoup...")
    soup = BeautifulSoup(html, "html.parser")
    log.debug("Rendering HTML elements to ReportLab story...")
    story = render_html_elements(soup)
    log.debug(f"Story generated with {len(story)} elements.")
    if not story:
        log.warning("Story is empty, adding 'Empty Content' paragraph.")
        story = [Paragraph("Empty Content", styles["CustomNormal"])]

    log.debug(f"Creating SimpleDocTemplate at {filepath}...")
    doc = SimpleDocTemplate(
        filepath,
        topMargin=72,
        bottomMargin=72,
        leftMargin=72,
        rightMargin=72
    )
    try:
        log.debug("Attempting to build PDF document...")
        log.debug(f"Calling doc.build with story containing {len(story)} elements.")
        doc.build(story)
        log.debug(f"PDF creation succeed: {filepath}")
    except Exception as e:
        log.error(f"Error in PDF building: {e}", exc_info=True) # Include traceback
        log.debug("Attempting to build PDF with error message...")
        simple_story = [Paragraph("Error in PDF generation", styles["CustomNormal"])]
        try:
            doc.build(simple_story)
            log.debug("Error PDF created successfully.")
        except Exception as e2:
            log.error(f"Failed to create even the error PDF: {e2}", exc_info=True)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    log.debug("create_pdf tool finished.")
    return {"url": _public_url(folder_path, fname)}

@mcp.tool()
def create_file(content: str, filename: str, persistent: bool = PERSISTENT_FILES) -> dict:
    folder_path = _generate_unique_folder()
    base, ext = os.path.splitext(filename)
    filepath = os.path.join(folder_path, filename)
    counter = 1
    while os.path.exists(filepath):
        filename = f"{base}_{counter}{ext}"
        filepath = os.path.join(folder_path, filename)
        counter += 1
    if ext.lower() == ".xml" and not content.strip().startswith("<?xml"):
        content = f'<?xml version="1.0" encoding="UTF-8"?>\n{content}'
    with open(filepath, "w", encoding="utf-8") as f:
        f.write(content)
    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    return {"url": _public_url(folder_path, filename)}

@mcp.tool()
def create_presentation(slides_data: list[dict], filename: str = None, persistent: bool = PERSISTENT_FILES, title: str = None) -> dict:
    folder_path = _generate_unique_folder()
    filepath, fname = _generate_filename(folder_path, "pptx", filename)

    # --- normalize presentation + layouts (template or blank; single-slide template => use same layout for all) ---
    use_template = False
    prs = None
    title_layout = None
    content_layout = None

    if PPTX_TEMPLATE:
        try:
            src = PPTX_TEMPLATE
            if hasattr(PPTX_TEMPLATE, "slides") and hasattr(PPTX_TEMPLATE, "save"):
                buf = BytesIO()
                PPTX_TEMPLATE.save(buf); buf.seek(0)
                src = buf

            tmp = Presentation(src)
            if len(tmp.slides) >= 1:
                prs = tmp
                use_template = True

                # If the template has 2+ slides: slide 0 = title layout, slide 1 = content layout
                # If it has exactly 1 slide: use slide 0 layout for BOTH title and content
                title_layout = prs.slides[0].slide_layout
                content_layout = prs.slides[1].slide_layout if len(prs.slides) >= 2 else prs.slides[0].slide_layout

                # Keep only the first slide (as title base)
                for i in range(len(prs.slides) - 1, 0, -1):
                    rId = prs.slides._sldIdLst[i].rId  # type: ignore[attr-defined]
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]        # type: ignore[attr-defined]
            # else -> fall through to no-template
        except Exception:
            use_template = False
            prs = None

    if not use_template:
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

    # Title slide (either existing template title, or newly added)
    if use_template:
        tslide = prs.slides[0]
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28); r.font.bold = True
    else:
        tslide = prs.slides.add_slide(title_layout)
        if tslide.shapes.title:
            tslide.shapes.title.text = title or ""
            for p in tslide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28); r.font.bold = True

    # slide size in inches (for robust layout math)
    EMU_PER_IN = 914400
    slide_w_in = prs.slide_width / EMU_PER_IN
    slide_h_in = prs.slide_height / EMU_PER_IN

    # shared margins/gutters
    page_margin = 0.5   # outer margin on each side (inches)
    gutter = 0.3        # space between image and text (inches)

    # --- shared path: add content slides ---
    for slide_data in slides_data:
        if not isinstance(slide_data, dict):
            raise ValueError("Each slide must be a dictionary.")

        slide_title = slide_data.get("title", "Untitled")
        content_list = slide_data.get("content", [])
        if not isinstance(content_list, list):
            content_list = [content_list]

        slide = prs.slides.add_slide(content_layout)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_title
            for p in slide.shapes.title.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(28); r.font.bold = True

        # Find or create a content shape
        content_shape = None
        try:
            for ph in slide.placeholders:
                try:
                    if ph.placeholder_format.idx == 1:
                        content_shape = ph; break
                except Exception:
                    pass
            if content_shape is None:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.idx != 0:
                            content_shape = ph; break
                    except Exception:
                        pass
        except Exception:
            pass
        if content_shape is None:
            content_shape = slide.shapes.add_textbox(Inches(page_margin), Inches(1.5), Inches(slide_w_in - 2*page_margin), Inches(slide_h_in - 2.0))

        # prep text frame: wrap + shrink-to-fit + small inner margins
        tf = content_shape.text_frame
        try:
            tf.clear()
        except Exception:
            pass
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        try:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05)
            tf.margin_bottom = Inches(0.05)
        except Exception:
            pass

        # default content box (will adjust if image present)
        content_left_in, content_top_in = page_margin, 1.5
        content_width_in = slide_w_in - 2*page_margin
        content_height_in = slide_h_in - (1.5 + page_margin)

        # Optional image placement with proper content reflow
        image_query = slide_data.get("image_query")
        if image_query:
            image_url = search_image(image_query)
            if image_url:
                try:
                    img = requests.get(image_url, timeout=10).content
                    stream = BytesIO(img)
                    pos = (slide_data.get("image_position") or "right").lower()
                    size = (slide_data.get("image_size") or "medium").lower()
                    if size == "small":
                        img_w_in, img_h_in = 2.0, 1.5
                    elif size == "large":
                        img_w_in, img_h_in = 4.0, 3.0
                    else:
                        img_w_in, img_h_in = 3.0, 2.0

                    if pos == "left":
                        img_left_in = page_margin
                        img_top_in = 1.5
                        content_left_in = img_left_in + img_w_in + gutter
                        content_top_in = 1.5
                        content_width_in = max(slide_w_in - page_margin - content_left_in, 2.5)
                        content_height_in = slide_h_in - (1.5 + page_margin)
                    elif pos == "right":
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = 1.5
                        content_left_in = page_margin
                        content_top_in = 1.5
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (1.5 + page_margin)
                    elif pos == "top":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = page_margin
                        content_left_in = page_margin
                        content_top_in = img_top_in + img_h_in + gutter
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(slide_h_in - page_margin - content_top_in, 2.0)
                    elif pos == "bottom":
                        img_left_in = slide_w_in - page_margin - img_w_in
                        img_top_in = max(slide_h_in - page_margin - img_h_in, page_margin)
                        content_left_in = page_margin
                        content_top_in = 1.0
                        content_width_in = slide_w_in - 2*page_margin
                        content_height_in = max(img_top_in - gutter - content_top_in, 2.0)
                    else:
                        img_left_in = max(slide_w_in - page_margin - img_w_in, page_margin)
                        img_top_in = 1.5
                        content_left_in = page_margin
                        content_top_in = 1.5
                        content_width_in = max(img_left_in - gutter - content_left_in, 2.5)
                        content_height_in = slide_h_in - (1.5 + page_margin)

                    slide.shapes.add_picture(stream, Inches(img_left_in), Inches(img_top_in), Inches(img_w_in), Inches(img_h_in))
                except Exception:
                    pass

        # apply content box geometry
        try:
            content_shape.left = Inches(content_left_in)
            content_shape.top = Inches(content_top_in)
            content_shape.width = Inches(content_width_in)
            content_shape.height = Inches(content_height_in)
        except Exception:
            pass

        # estimate capacity to guide initial font size; autosize will fine-tune
        approx_chars_per_in = 9.5
        approx_lines_per_in = 1.6
        est_capacity = int(content_width_in * approx_chars_per_in * content_height_in * approx_lines_per_in)
        font_size = dynamic_font_size(content_list, max_chars=max(est_capacity, 120), base_size=24, min_size=12)

        if not tf.paragraphs:
            tf.add_paragraph()
        for idx, line in enumerate(content_list):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = font_size
            run.font.name = "Calibri"
            p.space_after = Pt(6)

    prs.save(filepath)
    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    return {"url": _public_url(folder_path, fname)}

@mcp.tool()
def generate_and_archive(files_data: list[dict], archive_format: str = "zip", archive_name: str = None, persistent: bool = PERSISTENT_FILES) -> dict:
    folder_path = _generate_unique_folder()
    generated_files = []
    for file_info in files_data:
        filename = file_info.get("filename")
        content = file_info.get("content")
        format_type = file_info.get("format")
        title_param = file_info.get("title") 
        if content is None:
            content = ""    
        if title_param is None:
            title_param = ""
        filepath = os.path.join(folder_path, filename)
        os.makedirs(os.path.dirname(filepath), exist_ok=True)  
        try:
            if format_type == "py" or format_type == "cs" or format_type == "txt":
                with open(filepath, "w", encoding="utf-8") as f:
                    f.write(content)
            elif format_type == "pdf":
                if isinstance(content, list):
                    md_text = "\n".join(content)
                else:
                    md_text = content

                def replace_image_query(match):
                    query = match.group(1).strip()
                    log.debug(f"Found image_query placeholder: '{query}'")
                    image_url = search_image(query)
                    if image_url:
                        tag = f'<img src="{image_url}" alt="Image search: {query}" />'
                        log.debug(f"Replaced image_query '{query}' with {image_url}")
                    else:
                        tag = f'<img src="" alt="Image not found for: {query}" />'
                        log.warning(f"No image found for '{query}'")
                    return tag

                md_text = re.sub(r'!\[[^\]]*\]\(image_query:([^)]+)\)', replace_image_query, md_text)
                log.debug(f"Markdown after replacement for {filename}:\n{md_text}")

                html = markdown2.markdown(
                    md_text,
                    extras=[
                        'fenced-code-blocks',
                        'tables',
                        'break-on-newline',
                        'cuddled-lists',
                        'smarty-pants'
                    ]
                )
                log.debug(f"HTML generated for {filename}:\n{html}")

                soup = BeautifulSoup(html, "html.parser")
                story = render_html_elements(soup)

                if not story:
                    log.warning(f"Story empty for {filename}, adding fallback text.")
                    story = [Paragraph("Empty Content", styles["CustomNormal"])]

                doc = SimpleDocTemplate(
                    filepath,
                    topMargin=72,
                    bottomMargin=72,
                    leftMargin=72,
                    rightMargin=72
                )
                try:
                    doc.build(story)
                    log.debug(f"PDF '{filename}' successfully created in the archive.")
                except Exception as e:
                    log.error(f"Error during PDF build for '{filename}': {e}", exc_info=True)
                    fallback_story = [Paragraph("Error generating PDF", styles["CustomNormal"])]
                    doc.build(fallback_story)
            elif format_type == "xlsx":
                wb = Workbook()
                ws = wb.active
                if isinstance(content, list):
                    for row in content:
                        ws.append(row)
                wb.save(filepath)
            elif format_type == "csv":
                with open(filepath, "w", newline="", encoding="utf-8") as f:
                    if isinstance(content, list):
                        csv.writer(f).writerows(content)
                    else:
                        csv.writer(f).writerow([content])
            elif format_type == "pptx":
                parsed_content = file_info.get("slides_data", content)
                if isinstance(parsed_content, str):
                    try:
                        parsed_content = ast.literal_eval(parsed_content)
                        if not isinstance(parsed_content, list):
                            raise ValueError("Parsed content is not a list")
                    except (ValueError, SyntaxError):
                        raise ValueError(
                            f"Invalid format for pptx content: expected list of dicts, got '{type(parsed_content).__name__}'"
                        )
                elif not isinstance(parsed_content, list):
                    raise ValueError(f"Invalid format for pptx content: expected list, got '{type(parsed_content).__name__}'")
                prs = Presentation()
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title_shape = slide.shapes.title
                title_shape.text = title_param or "Presentation"
                for slide_data in parsed_content:
                    if not isinstance(slide_data, dict):
                        raise ValueError("Each slide must be a dictionary.")
                    title = slide_data.get("title", "Untitled")
                    content_list = slide_data.get("content", [])
                    if not isinstance(content_list, list):
                        content_list = [content_list]
                        
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title_shape = slide.shapes.title
                    title_shape.text = title
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(28)
                            run.font.bold = True
                            
                    content_shape = slide.placeholders[1]
                    content_shape.text = ""
                    font_size = dynamic_font_size(content_list, max_chars=300, base_size=24, min_size=12)

                    for line in content_list:
                        p = content_shape.text_frame.add_paragraph()
                        run = p.add_run()
                        run.text = line
                        run.font.size = font_size
                        p.space_after = Pt(6)
                        
                    image_query = slide_data.get("image_query")
                    if image_query:
                        image_url = search_image(image_query)
                        if image_url:
                            image_data = requests.get(image_url).content
                            image_stream = BytesIO(image_data)
                            position = slide_data.get("image_position", "right")
                            size = slide_data.get("image_size", "medium")
                            if size == "small":
                                width = Inches(2)
                                height = Inches(1.5)
                            elif size == "large":
                                width = Inches(4)
                                height = Inches(3)
                            else:
                                width = Inches(3)
                                height = Inches(2)
                            if position == "left":
                                left = Inches(0.5)
                                top = Inches(1.5)
                                content_shape.left = Inches(4.5)
                                content_shape.top = Inches(1.5)
                                content_shape.width = Inches(5)
                                content_shape.height = Inches(4)
                            elif position == "right":
                                left = Inches(5.5)
                                top = Inches(1.5)
                                content_shape.left = Inches(0.5)
                                content_shape.top = Inches(1.5)
                                content_shape.width = Inches(5)
                                content_shape.height = Inches(4)
                            elif position == "top":
                                left = Inches(5.5)
                                top = Inches(0.5)
                                content_shape.left = Inches(0.5)
                                content_shape.top = Inches(2.5)
                                content_shape.width = Inches(7)
                                content_shape.height = Inches(3)
                            elif position == "bottom":
                                left = Inches(5.5)
                                top = Inches(4.5)
                                content_shape.left = Inches(0.5)
                                content_shape.top = Inches(0.5)
                                content_shape.width = Inches(7)
                                content_shape.height = Inches(3)
                            else:
                                left = Inches(5.5)
                                top = Inches(1.5)
                                content_shape.left = Inches(0.5)
                                content_shape.top = Inches(1.5)
                                content_shape.width = Inches(5)
                                content_shape.height = Inches(4)
                            slide.shapes.add_picture(image_stream, left, top, width, height)
                    else:
                        content_shape.left = Inches(0.5)
                        content_shape.top = Inches(1.5)
                        content_shape.width = Inches(7)
                        content_shape.height = Inches(4)
                prs.save(filepath)
            else:
                with open(filepath, "w", encoding="utf-8") as f:
                    f.write(content)
            generated_files.append(filepath)
        except Exception as e:
            log.error(f"Error processing file '{filename}': {e}")
            raise 
    timestamp = datetime.datetime.now().strftime('%Y%m%d_%H%M%S')
    if archive_format.lower() == "7z":
        archive_filename = f"{archive_name or 'archive'}_{timestamp}.7z"
        archive_path = os.path.join(folder_path, archive_filename)
        with py7zr.SevenZipFile(archive_path, mode='w') as archive:
            for file_path in generated_files:
                archive.write(file_path, os.path.relpath(file_path, folder_path))
    elif archive_format.lower() == "tar.gz":
        archive_filename = f"{archive_name or 'archive'}_{timestamp}.tar.gz"
        archive_path = os.path.join(folder_path, archive_filename)
        with tarfile.open(archive_path, "w:gz") as tar:
            for file_path in generated_files:
                tar.add(file_path, arcname=os.path.relpath(file_path, folder_path))
    else: 
        archive_filename = f"{archive_name or 'archive'}_{timestamp}.zip"
        archive_path = os.path.join(folder_path, archive_filename)
        with zipfile.ZipFile(archive_path, 'w') as zipf:
            for file_path in generated_files:
                zipf.write(file_path, os.path.relpath(file_path, folder_path))
    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)
    return {"url": _public_url(folder_path, archive_filename)}

if __name__ == "__main__":
    mcp.run()