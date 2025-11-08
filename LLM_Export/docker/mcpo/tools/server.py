import os
import re
import json
import uuid
import shutil
import tarfile
import zipfile
import logging
from io import BytesIO
from typing import Any, List, Optional, Tuple

import py7zr
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.session import ServerSession

from docx import Document
from docx.shared import Inches
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt



from utils import (
    # generators
    create_presentation,
    create_word,
    create_excel,
    create_pdf,
    # file ops
    _create_csv,
    _create_raw_file,
    upload_file,
    download_file,
    _public_url,
    _generate_unique_folder,
    _generate_filename,
    _cleanup_files,
    # search / misc
    search_image,
    # xlsx helpers
    add_auto_sized_review_comment,
    # pptx helpers
    ensure_slot_textbox,
    _set_text_with_runs,
    _add_table_from_matrix,
    _set_table_from_matrix,
    shape_by_id,
    _collect_needs,
    _pick_layout_for_slots,
    _get_pptx_namespaces,
    _add_native_pptx_comment_zip,
    # docx helpers
    _apply_text_to_paragraph,
    _apply_run_formatting,
    _extract_paragraph_style_info,
)
from utils.pptx_treatment import _resolve_donor_simple

SCRIPT_VERSION = "0.9.0-fc1"

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv("LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s")

def _resolve_log_level(val: Optional[str]) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(level=_resolve_log_level(LOG_LEVEL_ENV), format=LOG_FORMAT_ENV)
log = logging.getLogger("server")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

URL = os.getenv('OWUI_URL')
TOKEN = os.getenv('JWT_SECRET')

PERSISTENT_FILES = os.getenv("PERSISTENT_FILES", "false")
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)

BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

DOCS_TEMPLATE_PATH = os.getenv("DOCS_TEMPLATE_DIR", "/rootPath/templates")
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        try:
            PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
            logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"PPTX template failed to load : {e}")
            PPTX_TEMPLATE = None
    else:
        logging.debug("No PPTX template found. Creation of a blank document.")
        PPTX_TEMPLATE = None

    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    XLSX_TEMPLATE_PATH = os.path.join("/rootPath/templates","Default_Template.xlsx")

    if XLSX_TEMPLATE_PATH:
        try:
            XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
            logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"Failed to load XLSX template: {e}")
            XLSX_TEMPLATE = None
    else:
        logging.debug("No XLSX template found. Creation of a blank document.")
        XLSX_TEMPLATE = None

# -----------------------------------------------------------------------------
# MCP server
# -----------------------------------------------------------------------------

mcp = FastMCP("file_export")


@mcp.tool(
    name="full_context_document",
    title="Return the structure of a document (docx, xlsx, pptx)",
    description="Return the structure, content, and metadata of a document based on its type (docx, xlsx, pptx). Unified output format with index, type, style, and text."
)
async def full_context_document(
    file_id: str,
    file_name: str,
    mcpo_headers: dict | None = None,
    ctx: Context[ServerSession, None] | None = None
) -> dict:
    """
    Inspect a document structure (docx/xlsx/pptx) and return a unified JSON representation.
    """
    user_token = TOKEN
    if mcpo_headers:
        auth_header = mcpo_headers.get("authorization")
        if auth_header:
            user_token = auth_header
            log.info("Using authorization from MCPO forwarded headers")
        else:
            log.warning("Forwarded headers present but no authorization found")
    else:
        log.info("No forwarded headers, using admin TOKEN fallback")

    try:
        user_file = download_file(file_id=file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        structure: dict[str, Any] = {
            "file_name": file_name,
            "file_id": file_id,
            "type": file_type,
            "slide_id_order": [],
            "body": [],
        }

        if file_type == "docx":
            doc = Document(user_file)
            para_id_counter = 1
            for para in doc.paragraphs:
                text = (para.text or "").strip()
                if not text:
                    continue
                style = getattr(para.style, "name", "")
                style_info = _extract_paragraph_style_info(para)
                element_type = "heading" if style.startswith("Heading") else "paragraph"
                para_xml_id = para_id_counter
                structure["body"].append({
                    "index": para_id_counter,
                    "para_xml_id": para_xml_id,
                    "id_key": f"pid:{para_xml_id}",
                    "type": element_type,
                    "style": style,
                    "style_info": style_info,
                    "text": text
                })
                para_id_counter += 1

            for table_idx, table in enumerate(doc.tables):
                table_xml_id = id(table._element)
                table_info = {
                    "index": para_id_counter,
                    "table_xml_id": table_xml_id,
                    "id_key": f"tid:{table_xml_id}",
                    "type": "table",
                    "style": "Table",
                    "table_id": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "cells": []
                }
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_xml_id = id(cell._element)
                        cell_text = (cell.text or "").strip()
                        cell_data = {
                            "row": row_idx,
                            "column": col_idx,
                            "cell_xml_id": cell_xml_id,
                            "id_key": f"tid:{table_xml_id}/cid:{cell_xml_id}",
                            "text": cell_text,
                            "style": getattr(getattr(cell, "style", None), "name", None)
                        }
                        row_data.append(cell_data)
                    table_info["cells"].append(row_data)
                structure["body"].append(table_info)
                para_id_counter += 1

        elif file_type == "xlsx":
            wb = load_workbook(user_file, read_only=True, data_only=True)
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    for col_idx, cell in enumerate(row, start=1):
                        if cell is None or str(cell).strip() == "":
                            continue
                        col_letter = sheet.cell(row=row_idx, column=col_idx).column_letter
                        cell_ref = f"{col_letter}{row_idx}"
                        structure["body"].append({
                            "index": cell_ref,
                            "type": "cell",
                            "text": str(cell)
                        })

        elif file_type == "pptx":
            prs = Presentation(user_file)
            structure["slide_id_order"] = [int(s.slide_id) for s in prs.slides]
            for slide_idx, slide in enumerate(prs.slides):
                title_shape = getattr(slide.shapes, "title", None)
                title_text = ""
                try:
                    title_text = (getattr(title_shape, "text", "") or "").strip() if title_shape else ""
                except Exception:
                    title_text = ""

                slide_obj: dict[str, Any] = {
                    "index": slide_idx,
                    "slide_id": int(slide.slide_id),
                    "id_key": f"sid:{int(slide.slide_id)}",
                    "title": title_text,
                    "shapes": []
                }

                for shape_idx, shape in enumerate(slide.shapes):
                    key = f"s{slide_idx}/sh{shape_idx}"
                    if hasattr(shape, "image"):
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "image"
                        })
                        continue

                    if hasattr(shape, "text_frame") and shape.text_frame:
                        kind = "title" if (title_shape is not None and shape is title_shape) else "textbox"
                        paragraphs = []
                        for p in shape.text_frame.paragraphs:
                            try:
                                text = "".join(r.text for r in p.runs) if p.runs else p.text
                            except Exception:
                                text = getattr(p, "text", "") or ""
                            text = (text or "").strip()
                            if text:
                                paragraphs.append(text)
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": kind,
                            "paragraphs": paragraphs
                        })
                        continue

                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        rows = []
                        for r in tbl.rows:
                            row_cells = []
                            for c in r.cells:
                                if hasattr(c, "text_frame") and c.text_frame:
                                    paras = []
                                    for p in c.text_frame.paragraphs:
                                        t = "".join(run.text for run in p.runs) if p.runs else p.text
                                        t = (t or "").strip()
                                        if t:
                                            paras.append(t)
                                    cell_text = "\n".join(paras)
                                else:
                                    cell_text = (getattr(c, "text", "") or "").strip()
                                row_cells.append(cell_text)
                            rows.append(row_cells)
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "table",
                            "rows": rows
                        })
                        continue

                structure["body"].append(slide_obj)

        else:
            return json.dumps({
                "error": {"message": f"Unsupported file type: {file_type}. Only docx, xlsx, and pptx are supported."}
            }, indent=4, ensure_ascii=False)

        return json.dumps(structure, indent=4, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)


@mcp.tool()
async def edit_document(
    file_id: str,
    file_name: str,
    edits: dict | list,
    mcpo_headers: dict | None = None,
    ctx: Context[ServerSession, None] | None = None
) -> dict:
    """
    Apply structural/content edits to a DOCX/XLSX/PPTX document.
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    user_token = TOKEN
    if mcpo_headers:
        auth_header = mcpo_headers.get("authorization")
        if auth_header:
            user_token = auth_header
            log.info("Using authorization from MCPO forwarded headers")
        else:
            log.warning("Forwarded headers present but no authorization found")
    else:
        log.info("No forwarded headers, using admin TOKEN fallback")

    try:
        user_file = download_file(file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        response: dict | None = None

        if file_type == "docx":
            try:
                doc = Document(user_file)

                para_by_xml_id: dict[int, Any] = {}
                table_by_xml_id: dict[int, Any] = {}
                cell_by_xml_id: dict[int, Any] = {}

                para_id_counter = 1
                for para in doc.paragraphs:
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for table in doc.tables:
                    table_xml_id = id(table._element)
                    table_by_xml_id[table_xml_id] = table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_xml_id = id(cell._element)
                            cell_by_xml_id[cell_xml_id] = cell

                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits

                new_refs: dict[str, int] = {}

                # ops: insert_after/before/delete_paragraph
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]

                    if kind == "insert_after" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            para_index = doc.paragraphs.index(anchor_para)
                            new_para = doc.add_paragraph()
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element) + 1, new_para._element)
                            new_para.style = anchor_para.style
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para

                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            new_para = doc.add_paragraph()
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element), new_para._element)
                            new_para.style = anchor_para.style
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para

                    elif kind == "delete_paragraph" and len(op) >= 2:
                        para_xml_id = int(op[1])
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            parent = para._element.getparent()
                            parent.remove(para._element)
                            para_by_xml_id.pop(para_xml_id, None)

                # content edits
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
                    t = target.strip()

                    m = re.match(r"^pid:(\d+)$", t, flags=re.I)
                    if m:
                        para_xml_id = int(m.group(1))
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            _apply_text_to_paragraph(para, new_text)
                        continue

                    m = re.match(r"^tid:(\d+)/cid:(\d+)$", t, flags=re.I)
                    if m:
                        table_xml_id = int(m.group(1))
                        cell_xml_id = int(m.group(2))
                        cell = cell_by_xml_id.get(cell_xml_id)
                        if cell:
                            for para in cell.paragraphs:
                                for _ in range(len(para.runs)):
                                    para._element.remove(para.runs[0]._element)
                            if cell.paragraphs:
                                first_para = cell.paragraphs[0]
                                first_para.add_run(str(new_text))
                        continue

                    m = re.match(r"^n(\d+)$", t, flags=re.I)
                    if m:
                        new_ref = t
                        para_xml_id = new_refs.get(new_ref)
                        if para_xml_id:
                            para = para_by_xml_id.get(para_xml_id)
                            if para:
                                _apply_text_to_paragraph(para, new_text)
                        continue

                edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.docx")
                doc.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="docx",
                    token=user_token,
                )

            except Exception as e:
                raise Exception(f"Error during DOCX editing: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active
                edit_items = edits.get("content_edits", []) if isinstance(edits, dict) and "content_edits" in edits else edits

                for index, new_text in edit_items:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"
                        cell = ws[cell_ref]
                        cell.value = new_text
                    except Exception:
                        ws["A1"].value = new_text

                edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.xlsx")
                wb.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="xlsx",
                    token=user_token,
                )

            except Exception as e:
                raise Exception(f"Error during XLSX editing: {e}")

        elif file_type == "pptx":
            try:
                prs = Presentation(user_file)

                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits

                new_ref_needs = _collect_needs(edit_items)
                order = [int(s.slide_id) for s in prs.slides]
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                new_refs: dict[str, int] = {}

                # ops: insert_before/after/delete_slide
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]

                    if kind in ("insert_after", "insert_before") and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            like_sid = None
                            if len(op) >= 4 and isinstance(op[3], dict):
                                like_sid = op[3].get("layout_like_sid")

                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})

                            if like_sid and like_sid in slides_by_id:
                                style_donor = slides_by_id[like_sid]
                            else:
                                style_donor = _resolve_donor_simple(order, slides_by_id, anchor_id, kind)

                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"]) if style_donor else prs.slide_layouts[0]
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                if kind == "insert_after":
                                    sldIdLst.insert(anchor_pos + 1, new_sldId)
                                    order.insert(anchor_pos + 1, new_sid)
                                else:
                                    sldIdLst.insert(anchor_pos, new_sldId)
                                    order.insert(anchor_pos, new_sid)
                            except Exception:
                                pass

                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid

                    elif kind == "delete_slide" and len(op) >= 2:
                        sid = int(op[1])
                        if sid in order:
                            i = order.index(sid)
                            sldIdLst = prs.slides._sldIdLst
                            rId = sldIdLst[i].rId
                            prs.part.drop_rel(rId)
                            del sldIdLst[i]
                            order.pop(i)
                            slides_by_id.pop(sid, None)

                # content edits
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
                    t = target.strip()

                    # table edits: sid:<sid>/shid:<shid> with 2D data
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if slide:
                            shape = shape_by_id(slide, shape_id)
                            if shape and getattr(shape, "has_table", False) and isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                                _set_table_from_matrix(shape, new_text)
                                continue

                    # text edits
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if not slide:
                            continue
                        shape = shape_by_id(slide, shape_id)
                        if not shape:
                            continue
                        _set_text_with_runs(shape, new_text)
                        continue

                    m = re.match(r"^(n\d+):slot:(title|body)$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        slot = m.group(2).lower()
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        shape = ensure_slot_textbox(slide, slot)
                        tf = getattr(shape, "text_frame", None)
                        if tf is None:
                            continue
                        if isinstance(new_text, list):
                            try:
                                tf.clear()
                            except Exception:
                                pass
                            tf.text = str(new_text[0]) if new_text else ""
                            for line in new_text[1:]:
                                p = tf.add_paragraph()
                                p.text = str(line)
                                try:
                                    p.level = getattr(tf.paragraphs[0], "level", 0)
                                except Exception:
                                    pass
                        else:
                            tf.text = str(new_text)
                        continue

                    # create a table on a new slide at slot:table
                    m = re.match(r"^(n\d+):slot:table$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        if isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                            _add_table_from_matrix(slide, new_text)
                        continue

                edited_path = os.path.join(temp_folder, f"{os.path.splitext(file_name)[0]}_edited.pptx")
                prs.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="pptx",
                    token=user_token,
                )

            except Exception as e:
                raise Exception(f"Error during PPTX editing: {e}")

        else:
            raise Exception(f"File type not supported: {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)
        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)


@mcp.tool(
    name="review_document",
    title="Review and comment on various document types",
    description="Review an existing document of various types (docx, xlsx, pptx), perform corrections and add comments."
)
async def review_document(
    file_id: str,
    file_name: str,
    review_comments: list[tuple[int | str, str]],
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Generic document review function that works with different document types.
    File type is automatically detected from the file extension.
    Returns a markdown hyperlink for downloading the reviewed document.
    
    For Excel files (.xlsx):
    - The index must be a cell reference (e.g. "A1", "B3", "C10")
    - These correspond to the "index" key returned by the full_context_document() function
    - Never use integer values for Excel cells
    
    For Word files (.docx):
    - The index should be a paragraph ID in the format "pid:<para_xml_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    
    For PowerPoint files (.pptx):
    - The index should be a slide ID in the format "sid:<slide_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)

    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header use admin fallback")
        user_token=TOKEN
    try:
        user_file = download_file(file_id=file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        reviewed_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                paragraphs = list(doc.paragraphs)
                para_by_xml_id = {}
                para_id_counter = 1
                
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for index, comment_text in review_comments:
                    if isinstance(index, int) and 0 <= index < len(paragraphs):
                        para = paragraphs[index]
                        if para.runs:
                            try:
                                doc.add_comment(
                                    runs=[para.runs[0]],
                                    text=comment_text,
                                    author="AI Reviewer",
                                    initials="AI"
                                )
                            except Exception:
                                para.add_run(f"  [AI Comment: {comment_text}]")
                    elif isinstance(index, str) and index.startswith("pid:"):
                        try:
                            para_xml_id = int(index.split(":")[1])
                            para = para_by_xml_id.get(para_xml_id)
                            if para and para.runs:
                                try:
                                    doc.add_comment(
                                        runs=[para.runs[0]],
                                        text=comment_text,
                                        author="AI Reviewer",
                                        initials="AI"
                                    )
                                except Exception:
                                    para.add_run(f"  [AI Comment: {comment_text}]")
                        except Exception:
                            if isinstance(index, int) and 0 <= index < len(paragraphs):
                                para = paragraphs[index]
                                if para.runs:
                                    try:
                                        doc.add_comment(
                                            runs=[para.runs[0]],
                                            text=comment_text,
                                            author="AI Reviewer",
                                            initials="AI"
                                        )
                                    except Exception:
                                        para.add_run(f"  [AI Comment: {comment_text}]")
                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.docx"
                )
                doc.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX revision: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                for index, comment_text in review_comments:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        add_auto_sized_review_comment(cell, comment_text, author="AI Reviewer")

                    except Exception:
                        fallback_cell = ws["A1"]
                        add_auto_sized_review_comment(fallback_cell, comment_text, author="AI Reviewer")

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.xlsx"
                )
                wb.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error: {e}")

        elif file_type == "pptx":
            try:
                temp_pptx = os.path.join(temp_folder, "temp_input.pptx")
                with open(temp_pptx, 'wb') as f:
                    f.write(user_file.read())
                
                prs = Presentation(temp_pptx)
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                
                comments_by_slide = {}
                
                for index, comment_text in review_comments:
                    slide_num = None
                    slide_id = None
                    
                    if isinstance(index, int) and 0 <= index < len(prs.slides):
                        slide_num = index + 1
                        slide_id = list(slides_by_id.keys())[index]
                    elif isinstance(index, str):
                        if index.startswith("sid:") and "/shid:" in index:
                            try:
                                slide_id = int(index.split("/")[0].replace("sid:", ""))
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse shape ID: {e}")
                        elif index.startswith("sid:"):
                            try:
                                slide_id = int(index.split(":")[1])
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse slide ID: {e}")
                    
                    if slide_num and slide_id:
                        if slide_num not in comments_by_slide:
                            comments_by_slide[slide_num] = []
                        
                        shape_info = ""
                        if "/shid:" in str(index):
                            try:
                                shape_id = int(str(index).split("/shid:")[1])
                                shape_info = f"[Shape {shape_id}] "
                            except:
                                pass
                        
                        comments_by_slide[slide_num].append(f"{shape_info}{comment_text}")
                comment_offset = 0              
                for slide_num, comments in comments_by_slide.items():
                    comment_start_x = 5000
                    comment_start_y = 1000
                    comment_spacing_y = 1500
                    
                    for idx, comment in enumerate(comments):
                        try:
                            y_position = comment_start_y + (idx * comment_spacing_y)
                            
                            _add_native_pptx_comment_zip(
                                pptx_path=temp_pptx,
                                slide_num=slide_num,
                                comment_text=f"• {comment}",
                                author_id=0,
                                x=comment_start_x,
                                y=y_position
                            )
                            log.debug(f"Native PowerPoint comment added to slide {slide_num} at position x={comment_start_x}, y={y_position}")
                        except Exception as e:
                            log.warning(f"Failed to add native comment to slide {slide_num}: {e}", exc_info=True)
                            prs_fallback = Presentation(temp_pptx)
                            slide = prs_fallback.slides[slide_num - 1]
                            left = top = Inches(0.2)
                            width = Inches(4)
                            height = Inches(1)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            p.text = f"AI Reviewer: {comment}"
                            p.font.size = PptPt(10)
                            prs_fallback.save(temp_pptx)

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.pptx"
                )
                shutil.copy(temp_pptx, reviewed_path)
                
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error when revising PPTX: {e}")

        else:
            raise Exception(f"File type not supported : {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )


@mcp.tool()
async def create_file(data: dict, persistent: bool = PERSISTENT_FILES) -> dict:
    """
    Create a single file based on 'data' description.
    data examples:
      {"format":"pdf","filename":"report.pdf","content":[...]}
      {"format":"docx","filename":"doc.docx","content":[...],"title":"..."}
      {"format":"pptx","filename":"slides.pptx","slides_data":[...],"title":"..."}
      {"format":"xlsx","filename":"data.xlsx","content":[[...]],"title":"..."}
      {"format":"csv","filename":"data.csv","content":[[...]]}
      {"format":"txt|xml|py|...","filename":"file.ext","content":"string"}
    """
    log.debug("Creating file via tool (server.py)")
    folder_path = _generate_unique_folder()
    format_type = (data.get("format") or "").lower()
    filename = data.get("filename")
    content = data.get("content")
    title = data.get("title")

    if format_type == "pdf":
        result = create_pdf(content if isinstance(content, list) else [str(content or "")], filename, folder_path=folder_path)
    elif format_type == "pptx":
        result = create_presentation(
            data.get("slides_data", []),
            filename,
            folder_path=folder_path,
            title=title,
            pptx_template_path=PPTX_TEMPLATE_PATH,
        )
    elif format_type == "docx":
        result = create_word(
            content if content is not None else [],
            filename,
            folder_path=folder_path,
            title=title,
            docx_template_path=DOCX_TEMPLATE_PATH,
        )
    elif format_type == "xlsx":
        result = create_excel(
            content if content is not None else [],
            filename,
            folder_path=folder_path,
            title=title,
            xlsx_template_path=XLSX_TEMPLATE_PATH if "xlsx_template_path" in create_excel.__code__.co_varnames else None,  # type: ignore
        )
    elif format_type == "csv":
        result = _create_csv(content if content is not None else [], filename, folder_path=folder_path)
    else:
        use_filename = filename or f"export.{format_type or 'txt'}"
        result = _create_raw_file(content if content is not None else "", use_filename, folder_path=folder_path)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}


@mcp.tool()
async def generate_and_archive(
    files_data: list[dict],
    archive_format: str = "zip",
    archive_name: str | None = None,
    persistent: bool = PERSISTENT_FILES
) -> dict:
    """
    Generate multiple files then archive them.
    files_data: list of 'data' dicts (same shape as for create_file)
    archive_format: zip | 7z | tar.gz
    """
    log.debug("Generating archive via tool (server.py)")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")
        try:
            if fmt == "pdf":
                res = create_pdf(content if isinstance(content, list) else [str(content or "")], fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = create_presentation(
                    file_info.get("slides_data", []),
                    fname,
                    folder_path=folder_path,
                    title=title,
                    pptx_template_path=PPTX_TEMPLATE_PATH,
                )
            elif fmt == "docx":
                res = create_word(
                    content if content is not None else [],
                    fname,
                    folder_path=folder_path,
                    title=title,
                    docx_template_path=DOCX_TEMPLATE_PATH,
                )
            elif fmt == "xlsx":
                res = create_excel(
                    content if content is not None else [],
                    fname,
                    folder_path=folder_path,
                    title=title,
                    xlsx_template_path=XLSX_TEMPLATE_PATH if "xlsx_template_path" in create_excel.__code__.co_varnames else None,  # type: ignore
                )
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise
        generated_paths.append(res["path"])

    timestamp = __import__("datetime").datetime.now().strftime("%Y%m%d_%H%M%S")
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    if archive_format.lower() == "7z":
        archive_filename = f"{archive_basename}.7z"
    elif archive_format.lower() == "tar.gz":
        archive_filename = f"{archive_basename}.tar.gz"
    else:
        archive_filename = f"{archive_basename}.zip"

    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode="w") as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, "w") as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}


if __name__ == "__main__":
    log.info(f"Starting MCPO File Export Server v{SCRIPT_VERSION}")
    mcp.run()
