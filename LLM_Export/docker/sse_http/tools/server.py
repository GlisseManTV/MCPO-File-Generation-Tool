import os
import re
import json
import uuid
import shutil
import tarfile
import zipfile
import logging
from io import BytesIO
from typing import Any, List, Optional, Tuple, Union
from typing_extensions import TypedDict
import py7zr
from mcp.server.fastmcp import FastMCP, Context
from mcp.server.session import ServerSession

from docx import Document
from docx.shared import Inches
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt as PptPt

from utils import (
    # generators
    create_presentation,
    create_word,
    create_excel,
    create_pdf,
    # file ops
    _create_csv,
    _create_raw_file,
    upload_file,
    download_file,
    _public_url,
    _generate_unique_folder,
    _generate_filename,
    _cleanup_files,
    # search / misc
    search_image,
    # xlsx helpers
    add_auto_sized_review_comment,
    # pptx helpers
    ensure_slot_textbox,
    _set_text_with_runs,
    _add_table_from_matrix,
    _set_table_from_matrix,
    shape_by_id,
    _collect_needs,
    _pick_layout_for_slots,
    _get_pptx_namespaces,
    _add_native_pptx_comment_zip,
    # docx helpers
    _apply_text_to_paragraph,
    _apply_run_formatting,
    _extract_paragraph_style_info,
)
from utils.pptx_treatment import _resolve_donor_simple
#NonDockerImport
import asyncio
import uvicorn
from typing import Any
from mcp.server.sse import SseServerTransport
from starlette.requests import Request
from starlette.applications import Starlette
from starlette.routing import Route, Mount
from starlette.responses import Response, JSONResponse, StreamingResponse

SCRIPT_VERSION = "1.0.0-dev1"

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv("LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s")

def _resolve_log_level(val: Optional[str]) -> int:
    if not val:
        return logging.INFO
    v = val.strip()
    if v.isdigit():
        try:
            return int(v)
        except ValueError:
            return logging.INFO
    return getattr(logging, v.upper(), logging.INFO)

logging.basicConfig(level=_resolve_log_level(LOG_LEVEL_ENV), format=LOG_FORMAT_ENV)
log = logging.getLogger("server")
log.setLevel(_resolve_log_level(LOG_LEVEL_ENV))
log.info("Effective LOG_LEVEL -> %s", logging.getLevelName(log.level))

class ReviewComment(TypedDict):
    index: Union[int, str]
    comment: str

URL = os.getenv('OWUI_URL')
TOKEN = os.getenv('JWT_SECRET')

def _env_bool(val: str | None) -> bool:
    return str(val).strip().lower() in ("1", "true", "yes", "y", "on") if val is not None else False

PERSISTENT_FILES = _env_bool(os.getenv("PERSISTENT_FILES", "false"))
FILES_DELAY = int(os.getenv("FILES_DELAY", 60)) 

EXPORT_DIR_ENV = os.getenv("FILE_EXPORT_DIR")
EXPORT_DIR = (EXPORT_DIR_ENV or r"/output").rstrip("/")
os.makedirs(EXPORT_DIR, exist_ok=True)

BASE_URL_ENV = os.getenv("FILE_EXPORT_BASE_URL")
BASE_URL = (BASE_URL_ENV or "http://localhost:9003/files").rstrip("/")

LOG_LEVEL_ENV = os.getenv("LOG_LEVEL")
LOG_FORMAT_ENV = os.getenv(
    "LOG_FORMAT", "%(asctime)s %(levelname)s %(name)s - %(message)s"
)

DOCS_TEMPLATE_PATH = os.getenv("DOCS_TEMPLATE_DIR", "/rootPath/templates")
PPTX_TEMPLATE = None
DOCX_TEMPLATE = None
XLSX_TEMPLATE = None
PPTX_TEMPLATE_PATH = None
DOCX_TEMPLATE_PATH = None
XLSX_TEMPLATE_PATH = None

if DOCS_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
    logging.debug(f"Template Folder: {DOCS_TEMPLATE_PATH}")
    for root, dirs, files in os.walk(DOCS_TEMPLATE_PATH):
        for file in files:
            fpath = os.path.join(root, file)
            if file.lower().endswith(".pptx") and PPTX_TEMPLATE_PATH is None:
                PPTX_TEMPLATE_PATH = fpath
                logging.debug(f"PPTX template: {PPTX_TEMPLATE_PATH}")
            elif file.lower().endswith(".docx") and DOCX_TEMPLATE_PATH is None:
                DOCX_TEMPLATE_PATH = fpath
            elif file.lower().endswith(".xlsx") and XLSX_TEMPLATE_PATH is None:
                XLSX_TEMPLATE_PATH = fpath
    if PPTX_TEMPLATE_PATH:
        try:
            PPTX_TEMPLATE = Presentation(PPTX_TEMPLATE_PATH)
            logging.debug(f"Using PPTX template: {PPTX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"PPTX template failed to load : {e}")
            PPTX_TEMPLATE = None
    else:
        logging.debug("No PPTX template found. Creation of a blank document.")
        PPTX_TEMPLATE = None

    if DOCX_TEMPLATE_PATH and os.path.exists(DOCS_TEMPLATE_PATH):
        try:
            DOCX_TEMPLATE = Document(DOCX_TEMPLATE_PATH)
            logging.debug(f"Using DOCX template: {DOCX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"DOCX template failed to load : {e}")
            DOCX_TEMPLATE = None
    else:
        logging.debug("No DOCX template found. Creation of a blank document.")
        DOCX_TEMPLATE = None
    
    XLSX_TEMPLATE_PATH = os.path.join("/rootPath/templates","Default_Template.xlsx")

    if XLSX_TEMPLATE_PATH:
        try:
            XLSX_TEMPLATE = load_workbook(XLSX_TEMPLATE_PATH)
            logging.debug(f"Using XLSX template: {XLSX_TEMPLATE_PATH}")
        except Exception as e:
            logging.warning(f"Failed to load XLSX template: {e}")
            XLSX_TEMPLATE = None
    else:
        logging.debug("No XLSX template found. Creation of a blank document.")
        XLSX_TEMPLATE = None

# -----------------------------------------------------------------------------
# MCP server
# -----------------------------------------------------------------------------

mcp = FastMCP(
    name = "file_export",
    port = int(os.getenv("MCP_HTTP_PORT", "9004")),
    host = (os.getenv("MCP_HTTP_HOST", "0.0.0.0"))
)

@mcp.tool(
    name="full_context_document",
    title="Return the structure of a document (docx, xlsx, pptx)",
    description="Return the structure, content, and metadata of a document based on its type (docx, xlsx, pptx). Unified output format with index, type, style, and text."
)

async def full_context_document(
    file_id: str,
    file_name: str,
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Return the structure of a document (docx, xlsx, pptx) based on its file extension.
    The function detects the file type and processes it accordingly.
    Returns:
        dict: A JSON object with the structure of the document.
    """
    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        user_token=bearer_token
        logging.info(f"Recieved authorization header!")        
    except:
        user_token=TOKEN
        logging.error(f"Error retrieving authorization header use admin fallback")
    try:
        user_file = download_file(file_id=file_id,token=user_token)

        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        structure = {
            "file_name": file_name,
            "file_id": file_id,
            "type": file_type,
            "slide_id_order": [],
            "body": [],
        }
        index_counter = 1

        if file_type == "docx":
            doc = Document(user_file)
            
            para_id_counter = 1
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                style = para.style.name
                style_info = _extract_paragraph_style_info(para)
                element_type = "heading" if style.startswith("Heading") else "paragraph"
                
                para_xml_id = para_id_counter
                
                structure["body"].append({
                    "index": para_id_counter,
                    "para_xml_id": para_xml_id,
                    "id_key": f"pid:{para_xml_id}",
                    "type": element_type,
                    "style": style,
                    "style_info": style_info,
                    "text": text
                })
                para_id_counter += 1
            
            for table_idx, table in enumerate(doc.tables):
                table_xml_id = id(table._element)
                table_info = {
                    "index": para_id_counter,
                    "table_xml_id": table_xml_id,
                    "id_key": f"tid:{table_xml_id}",
                    "type": "table",
                    "style": "Table",
                    "table_id": table_idx,
                    "rows": len(table.rows),
                    "columns": len(table.rows[0].cells) if table.rows else 0,
                    "cells": []
                }
                
                for row_idx, row in enumerate(table.rows):
                    row_data = []
                    for col_idx, cell in enumerate(row.cells):
                        cell_xml_id = id(cell._element)
                        cell_text = cell.text.strip()
                        cell_data = {
                            "row": row_idx,
                            "column": col_idx,
                            "cell_xml_id": cell_xml_id,
                            "id_key": f"tid:{table_xml_id}/cid:{cell_xml_id}",
                            "text": cell_text,
                            "style": cell.style.name if hasattr(cell, 'style') else None
                        }
                        row_data.append(cell_data)
                    table_info["cells"].append(row_data)
                
                structure["body"].append(table_info)
                para_id_counter += 1

        elif file_type == "xlsx":
            wb = load_workbook(user_file, read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    for col_idx, cell in enumerate(row, start=1):
                        if cell is None or str(cell).strip() == "":
                            continue
                        col_letter = sheet.cell(row=row_idx, column=col_idx).column_letter
                        cell_ref = f"{col_letter}{row_idx}"
                        structure["body"].append({
                            "index": cell_ref,
                            "type": "cell",
                            "text": str(cell)
                        })
                        index_counter += 1

        elif file_type == "pptx":
            prs = Presentation(user_file)
            structure["slide_id_order"] = [int(s.slide_id) for s in prs.slides]
            for slide_idx, slide in enumerate(prs.slides):
                title_shape = slide.shapes.title if hasattr(slide.shapes, "title") else None
                title_text = title_shape.text.strip() if (title_shape and getattr(title_shape, "text", "").strip()) else ""

                slide_obj = {
                    "index": slide_idx,
                    "slide_id": int(slide.slide_id),
                    "id_key": f"sid:{int(slide.slide_id)}",
                    "title": title_text,
                    "shapes": []
                }
                

                for shape_idx, shape in enumerate(slide.shapes):
                    key = f"s{slide_idx}/sh{shape_idx}"
                    if hasattr(shape, "image"):
                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id  
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "image"
                        })
                        continue

                    if hasattr(shape, "text_frame") and shape.text_frame:
                        kind = "title" if (title_shape is not None and shape is title_shape) else "textbox"

                        paragraphs = []
                        for p in shape.text_frame.paragraphs:
                            text = "".join(run.text for run in p.runs) if p.runs else p.text
                            text = (text or "").strip()
                            if text != "":
                                paragraphs.append(text)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id 
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key, 
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": kind,
                            "paragraphs": paragraphs
                        })
                        continue
                    if getattr(shape, "has_table", False):
                        tbl = shape.table
                        rows = []
                        for r in tbl.rows:
                            row_cells = []
                            for c in r.cells:
                                # collect full text
                                if hasattr(c, "text_frame") and c.text_frame:
                                    paras = []
                                    for p in c.text_frame.paragraphs:
                                        t = "".join(run.text for run in p.runs) if p.runs else p.text
                                        t = (t or "").strip()
                                        if t:
                                            paras.append(t)
                                    cell_text = "\n".join(paras)
                                else:
                                    cell_text = (getattr(c, "text", "") or "").strip()
                                row_cells.append(cell_text)
                            rows.append(row_cells)

                        shape_id_val = getattr(shape, "shape_id", None) or shape._element.cNvPr.id
                        slide_obj["shapes"].append({
                            "shape_idx": shape_idx,
                            "shape_id": shape_id_val,
                            "idx_key": key,
                            "id_key": f"sid:{int(slide.slide_id)}/shid:{int(shape_id_val)}",
                            "kind": "table",
                            "rows": rows  # list of lists: each inner list = one row's cell texts
                        })
                        continue

                structure["body"].append(slide_obj)

        else:
            return json.dumps({
                "error": {"message": f"Unsupported file type: {file_type}. Only docx, xlsx, and pptx are supported."}
            }, indent=4, ensure_ascii=False)

        return json.dumps(structure, indent=4, ensure_ascii=False)

    except Exception as e:
        return json.dumps({"error": {"message": str(e)}}, indent=4, ensure_ascii=False)

@mcp.tool()
async def edit_document(
    file_id: str,
    file_name: str,
    edits: dict,
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Edits a document (docx, xlsx, pptx) using structured operations.

    Args:
        file_id: Unique identifier for the document.
        file_name: Name of the document file.
        edits: Dictionary with:
            - "ops": List of structural changes.
            - "content_edits": List of content updates.

    ## Supported Formats

    ### PPTX (PowerPoint)
    - ops: 
        - ["insert_after", <slide_id>, "nK", {"layout_like_sid": <slide_id>}]
        - ["insert_after", <slide_id>, "nK", {"layout_like_sid": <slide_id>}]
        - ["delete_slide", slide_id]
    - content_edits:
        - Edit a text shape
            ["sid:<slide_id>/shid:<shape_id>", text_or_list]
        - Edit a table
            ["sid:<slide_id>/shid:<shape_id>", [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]]
        - Edit title or body or table of a newly inserted slide
            ["nK:slot:title", text_or_list]
            ["nK:slot:body", text_or_list]
            ["nK:slot:table", [[row1_col1, row1_col2], [row2_col1, row2_col2], ...]]


    ### DOCX (Word)
    - ops:
        - ["insert_after", para_xml_id, "nK"]
        - ["insert_before", para_xml_id, "nK"]
        - ["delete_paragraph", para_xml_id]
    - content_edits:
        - ["pid:<para_xml_id>", text_or_list]
        - ["tid:<table_xml_id>/cid:<cell_xml_id>", text]
        - ["nK", text_or_list]

    ### XLSX (Excel)
    - ops:
        - ["insert_row", "sheet_name", row_idx]
        - ["delete_row", "sheet_name", row_idx]
        - ["insert_column", "sheet_name", col_idx]
        - ["delete_column", "sheet_name", col_idx]
    - content_edits:
        - ["<ref>", value]

    ## Notes
    - Always call `full_context_document()` first to get IDs.
    - Use cell refs like "A1", "B5".
    - Formatting is preserved.
    - Returns a download link to the edited file.
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)
    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header use admin fallback")
        user_token=TOKEN
    try:
        user_file = download_file(file_id=file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        edited_path = None
        response = None

        if file_type == "docx":
            try:
                doc = Document(user_file)
                
                para_by_xml_id = {}
                table_by_xml_id = {}
                cell_by_xml_id = {}
                
                para_id_counter = 1
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1
          
                for table in doc.tables:
                    table_xml_id = id(table._element)
                    table_by_xml_id[table_xml_id] = table
                    for row in table.rows:
                        for cell in row.cells:
                            cell_xml_id = id(cell._element)
                            cell_by_xml_id[cell_xml_id] = cell
          
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
                try:
                    if isinstance(edit_items, list) and (len(edit_items) == 0 or isinstance(edit_items[0], dict)):
                        edit_items = [
                            [item.get("target"), item.get("value")]
                            for item in (edit_items or [])
                            if isinstance(item, dict) and "target" in item and "value" in item
                        ]
                except Exception:
                    pass 

                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]
                    
                    if kind == "insert_after" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            para_index = doc.paragraphs.index(anchor_para)
          	       	      
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element) + 1, new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_xml_id = int(op[1])
                        new_ref = op[2]
                        
                        anchor_para = para_by_xml_id.get(anchor_xml_id)
                        if anchor_para:
                            new_para = doc.add_paragraph()
          	       	      
                            anchor_element = anchor_para._element
                            parent = anchor_element.getparent()
                            parent.insert(parent.index(anchor_element), new_para._element)
          	       	      
                            new_para.style = anchor_para.style
      	      				
                            new_xml_id = id(new_para._element)
                            new_refs[new_ref] = new_xml_id
                            para_by_xml_id[new_xml_id] = new_para
      	      		
                    elif kind == "delete_paragraph" and len(op) >= 2:
                        para_xml_id = int(op[1])
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            parent = para._element.getparent()
                            parent.remove(para._element)
                            para_by_xml_id.pop(para_xml_id, None)
                
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
      	      		
                    t = target.strip()
                    
                    m = re.match(r"^pid:(\d+)$", t, flags=re.I)
                    if m:
                        para_xml_id = int(m.group(1))
                        para = para_by_xml_id.get(para_xml_id)
                        if para:
                            _apply_text_to_paragraph(para, new_text)
                        continue
      	      		
                    m = re.match(r"^tid:(\d+)/cid:(\d+)$", t, flags=re.I)
                    if m:
                        table_xml_id = int(m.group(1))
                        cell_xml_id = int(m.group(2))
                        cell = cell_by_xml_id.get(cell_xml_id)
                        if cell:
                            for para in cell.paragraphs:
                                for _ in range(len(para.runs)):
                                    para._element.remove(para.runs[0]._element)
          	       	      
                            if cell.paragraphs:
                                first_para = cell.paragraphs[0]
                                first_para.add_run(str(new_text))
                        continue
      	      		
                    m = re.match(r"^n(\d+)$", t, flags=re.I)
                    if m:
                        new_ref = t
                        para_xml_id = new_refs.get(new_ref)
                        if para_xml_id:
                            para = para_by_xml_id.get(para_xml_id)
                            if para:
                                _apply_text_to_paragraph(para, new_text)
                        continue
          
                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.docx"
                )
                doc.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX editing: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                edit_items = edits.get("content_edits", []) if isinstance(edits, dict) and "content_edits" in edits else edits
                try:
                    if isinstance(edit_items, list) and (len(edit_items) == 0 or isinstance(edit_items[0], dict)):
                        edit_items = [
                            [item.get("target"), item.get("value")]
                            for item in (edit_items or [])
                            if isinstance(item, dict) and "target" in item and "value" in item
                        ]
                except Exception:
                    pass       
                for index, new_text in edit_items:
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        cell.value = new_text
                    except Exception:
                        fallback_cell = ws["A1"]
                        fallback_cell.value = new_text

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.xlsx"
                )
                wb.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during XLSX editing: {e}")

        elif file_type == "pptx":
            try:
                prs = Presentation(user_file)
                if isinstance(edits, dict):
                    ops = edits.get("ops", []) or []
                    edit_items = edits.get("content_edits", []) or []
                else:
                    ops = []
                    edit_items = edits
                try:
                    if isinstance(edit_items, list) and (len(edit_items) == 0 or isinstance(edit_items[0], dict)):
                        edit_items = [
                            [item.get("target"), item.get("value")]
                            for item in (edit_items or [])
                            if isinstance(item, dict) and "target" in item and "value" in item
                        ]
                except Exception:
                    pass                 
                new_ref_needs = _collect_needs(edit_items)
                order = [int(s.slide_id) for s in prs.slides]
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                new_refs = {}
                
                for op in ops:
                    if not isinstance(op, (list, tuple)) or not op:
                        continue
                    kind = op[0]

                    if kind == "insert_after" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            like_sid = None
                            if len(op) >= 4 and isinstance(op[3], dict):
                                like_sid = op[3].get("layout_like_sid")
                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})
                            # style_donor selection
                            if like_sid and like_sid in slides_by_id:
                                style_donor = slides_by_id[like_sid]
                            else:
                                style_donor = _resolve_donor_simple(order, slides_by_id, anchor_id, kind)
                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"]) if style_donor else prs.slide_layouts[0]
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos + 1, new_sldId)
                            except Exception:
                                pass 
                            order.insert(order.index(anchor_id) + 1, new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid


                    elif kind == "insert_before" and len(op) >= 3:
                        anchor_id = int(op[1])
                        new_ref = op[2]
                        if anchor_id in order:
                            like_sid = None
                            if len(op) >= 4 and isinstance(op[3], dict):
                                like_sid = op[3].get("layout_like_sid")

                            needs = new_ref_needs.get(new_ref, {"title": False, "body": False})

                            # style_donor selection
                            if like_sid and like_sid in slides_by_id:
                                style_donor = slides_by_id[like_sid]
                            else:
                                style_donor = _resolve_donor_simple(order, slides_by_id, anchor_id, kind)

                            layout = _pick_layout_for_slots(prs, style_donor, needs["title"], needs["body"]) if style_donor else prs.slide_layouts[0]
                            new_slide = prs.slides.add_slide(layout)
                            new_sid = int(new_slide.slide_id)

                            sldIdLst = prs.slides._sldIdLst
                            new_sldId = sldIdLst[-1]
                            try:
                                anchor_pos = order.index(anchor_id)
                                sldIdLst.remove(new_sldId)
                                sldIdLst.insert(anchor_pos, new_sldId)
                            except Exception:
                                pass

                            order.insert(order.index(anchor_id), new_sid)
                            slides_by_id[new_sid] = new_slide
                            new_refs[new_ref] = new_sid

                    elif kind == "delete_slide" and len(op) >= 2:
                        sid = int(op[1])
                        if sid in order:
                            i = order.index(sid)
                            sldIdLst = prs.slides._sldIdLst
                            rId = sldIdLst[i].rId
                            prs.part.drop_rel(rId)
                            del sldIdLst[i]
                            order.pop(i)
                            slides_by_id.pop(sid, None)   		 
      		
                for target, new_text in edit_items:
                    if not isinstance(target, str):
                        continue
                    t = target.strip()

                    # >>> ADD: table edit 
                    # target format: sid:<sid>/shid:<shid>  with new_text like [[row1...],[row2...],...]
                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if slide:
                            shape = shape_by_id(slide, shape_id)
                            if shape and getattr(shape, "has_table", False) and isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                                _set_table_from_matrix(shape, new_text)
                                continue
                            
                    # <<< END ADD

                    m = re.match(r"^sid:(\d+)/shid:(\d+)$", t, flags=re.I)
                    if m:
                        slide_id = int(m.group(1))
                        shape_id = int(m.group(2))
                        slide = slides_by_id.get(slide_id)
                        if not slide:
                            continue
                        shape = shape_by_id(slide, shape_id)
                        if not shape:
                            continue
                        _set_text_with_runs(shape, new_text)
                        continue
                    m = re.match(r"^(n\d+):slot:(title|body)$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        slot = m.group(2).lower()
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        shape = ensure_slot_textbox(slide, slot)
                        tf = getattr(shape, "text_frame", None)
                        if tf is None:
                            continue
                        if isinstance(new_text, list):
                            tf.clear()
                            tf.text = str(new_text[0]) if new_text else ""
                            for line in new_text[1:]:
                                p = tf.add_paragraph()
                                p.text = str(line)
                                try:
                                    p.level = getattr(tf.paragraphs[0], "level", 0)
                                except Exception:
                                    pass
                        else:
                            tf.text = str(new_text)
                        continue

                    # nK:table  (create a new table on a newly inserted slide nK)
                    m = re.match(r"^(n\d+):slot:table$", t, flags=re.I)
                    if m:
                        ref = m.group(1)
                        sid = new_refs.get(ref)
                        if not sid:
                            continue
                        slide = slides_by_id.get(sid)
                        if not slide:
                            continue
                        # new_text must be a 2D list
                        if isinstance(new_text, (list, tuple)) and new_text and isinstance(new_text[0], (list, tuple)):
                            _add_table_from_matrix(slide, new_text)
                        continue
 

                edited_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_edited.pptx"
                )
                prs.save(edited_path)
                response = upload_file(
                    file_path=edited_path,
                    filename=f"{os.path.splitext(file_name)[0]}_edited",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during PPTX editing: {e}")

        else:
            raise Exception(f"File type not supported: {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )

@mcp.tool(
    name="review_document",
    title="Review and comment on various document types",
    description="Review an existing document of various types (docx, xlsx, pptx), perform corrections and add comments."
)
async def review_document(
    file_id: str,
    file_name: str,
    review_comments: List[ReviewComment],
    ctx: Context[ServerSession, None]
) -> dict:
    """
    Generic document review function that works with different document types.
    File type is automatically detected from the file extension.
    Returns a markdown hyperlink for downloading the reviewed document.
    
    For Excel files (.xlsx):
    - The index must be a cell reference (e.g. "A1", "B3", "C10")
    - These correspond to the "index" key returned by the full_context_document() function
    - Never use integer values for Excel cells
    
    For Word files (.docx):
    - The index should be a paragraph ID in the format "pid:<para_xml_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    
    For PowerPoint files (.pptx):
    - The index should be a slide ID in the format "sid:<slide_id>"
    - These correspond to the "id_key" field returned by the full_context_document() function
    """
    temp_folder = f"/app/temp/{uuid.uuid4()}"
    os.makedirs(temp_folder, exist_ok=True)

    try:
        bearer_token = ctx.request_context.request.headers.get("authorization")
        logging.info(f"Recieved authorization header!")
        user_token=bearer_token
    except:
        logging.error(f"Error retrieving authorization header use admin fallback")
        user_token=TOKEN
    try:
        user_file = download_file(file_id=file_id, token=user_token)
        if isinstance(user_file, dict) and "error" in user_file:
            return json.dumps(user_file, indent=4, ensure_ascii=False)

        file_extension = os.path.splitext(file_name)[1].lower()
        file_type = file_extension.lstrip('.')

        reviewed_path = None
        response = None
        norm_comments: List[ReviewComment] = []
        try:
            if isinstance(review_comments, list):
                if len(review_comments) == 0:
                    norm_comments = []
                elif isinstance(review_comments[0], dict):
                    for item in review_comments:
                        if isinstance(item, dict) and "index" in item and "comment" in item:
                            norm_comments.append({"index": item["index"], "comment": str(item["comment"])})
                else:
                    # Legacy: [[index, comment], ...] or tuples
                    for item in review_comments:
                        if isinstance(item, (list, tuple)) and len(item) >= 2:
                            norm_comments.append({"index": item[0], "comment": str(item[1])})
        except Exception:
            pass
        if file_type == "docx":
            try:
                doc = Document(user_file)
                paragraphs = list(doc.paragraphs)
                para_by_xml_id = {}
                para_id_counter = 1
                
                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    para_by_xml_id[para_id_counter] = para
                    para_id_counter += 1

                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    if isinstance(index, int) and 0 <= index < len(paragraphs):
                        para = paragraphs[index]
                        if para.runs:
                            try:
                                doc.add_comment(
                                    runs=[para.runs[0]],
                                    text=comment_text,
                                    author="AI Reviewer",
                                    initials="AI"
                                )
                            except Exception:
                                para.add_run(f"  [AI Comment: {comment_text}]")
                    elif isinstance(index, str) and index.startswith("pid:"):
                        try:
                            para_xml_id = int(index.split(":")[1])
                            para = para_by_xml_id.get(para_xml_id)
                            if para and para.runs:
                                try:
                                    doc.add_comment(
                                        runs=[para.runs[0]],
                                        text=comment_text,
                                        author="AI Reviewer",
                                        initials="AI"
                                    )
                                except Exception:
                                    para.add_run(f"  [AI Comment: {comment_text}]")
                        except Exception:
                            if isinstance(index, int) and 0 <= index < len(paragraphs):
                                para = paragraphs[index]
                                if para.runs:
                                    try:
                                        doc.add_comment(
                                            runs=[para.runs[0]],
                                            text=comment_text,
                                            author="AI Reviewer",
                                            initials="AI"
                                        )
                                    except Exception:
                                        para.add_run(f"  [AI Comment: {comment_text}]")
                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.docx"
                )
                doc.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="docx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error during DOCX revision: {e}")

        elif file_type == "xlsx":
            try:
                wb = load_workbook(user_file)
                ws = wb.active

                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    try:
                        if isinstance(index, str) and re.match(r"^[A-Z]+[0-9]+$", index.strip().upper()):
                            cell_ref = index.strip().upper()
                        elif isinstance(index, int):
                            cell_ref = f"A{index+1}"
                        else:
                            cell_ref = "A1"

                        cell = ws[cell_ref]
                        add_auto_sized_review_comment(cell, comment_text, author="AI Reviewer")

                    except Exception:
                        fallback_cell = ws["A1"]
                        add_auto_sized_review_comment(fallback_cell, comment_text, author="AI Reviewer")

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.xlsx"
                )
                wb.save(reviewed_path)
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="xlsx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error: {e}")

        elif file_type == "pptx":
            try:
                temp_pptx = os.path.join(temp_folder, "temp_input.pptx")
                with open(temp_pptx, 'wb') as f:
                    f.write(user_file.read())
                
                prs = Presentation(temp_pptx)
                slides_by_id = {int(s.slide_id): s for s in prs.slides}
                
                comments_by_slide = {}
                
                for rc in norm_comments:
                    index = rc["index"]
                    comment_text = rc["comment"]
                    slide_num = None
                    slide_id = None
                    
                    if isinstance(index, int) and 0 <= index < len(prs.slides):
                        slide_num = index + 1
                        slide_id = list(slides_by_id.keys())[index]
                    elif isinstance(index, str):
                        if index.startswith("sid:") and "/shid:" in index:
                            try:
                                slide_id = int(index.split("/")[0].replace("sid:", ""))
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse shape ID: {e}")
                        elif index.startswith("sid:"):
                            try:
                                slide_id = int(index.split(":")[1])
                                if slide_id in slides_by_id:
                                    slide_num = list(slides_by_id.keys()).index(slide_id) + 1
                            except Exception as e:
                                log.warning(f"Failed to parse slide ID: {e}")
                    
                    if slide_num and slide_id:
                        if slide_num not in comments_by_slide:
                            comments_by_slide[slide_num] = []
                        
                        shape_info = ""
                        if "/shid:" in str(index):
                            try:
                                shape_id = int(str(index).split("/shid:")[1])
                                shape_info = f"[Shape {shape_id}] "
                            except:
                                pass
                        
                        comments_by_slide[slide_num].append(f"{shape_info}{comment_text}")
                comment_offset = 0              
                for slide_num, comments in comments_by_slide.items():
                    comment_start_x = 5000
                    comment_start_y = 1000
                    comment_spacing_y = 1500
                    
                    for idx, comment in enumerate(comments):
                        try:
                            y_position = comment_start_y + (idx * comment_spacing_y)
                            
                            _add_native_pptx_comment_zip(
                                pptx_path=temp_pptx,
                                slide_num=slide_num,
                                comment_text=f"• {comment}",
                                author_id=0,
                                x=comment_start_x,
                                y=y_position
                            )
                            log.debug(f"Native PowerPoint comment added to slide {slide_num} at position x={comment_start_x}, y={y_position}")
                        except Exception as e:
                            log.warning(f"Failed to add native comment to slide {slide_num}: {e}", exc_info=True)
                            prs_fallback = Presentation(temp_pptx)
                            slide = prs_fallback.slides[slide_num - 1]
                            left = top = Inches(0.2)
                            width = Inches(4)
                            height = Inches(1)
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            p = text_frame.add_paragraph()
                            p.text = f"AI Reviewer: {comment}"
                            p.font.size = PptPt(10)
                            prs_fallback.save(temp_pptx)

                reviewed_path = os.path.join(
                    temp_folder, f"{os.path.splitext(file_name)[0]}_reviewed.pptx"
                )
                shutil.copy(temp_pptx, reviewed_path)
                
                response = upload_file(
                    file_path=reviewed_path,
                    filename=f"{os.path.splitext(file_name)[0]}_reviewed",
                    file_type="pptx", 
                    token=user_token
                )
            except Exception as e:
                raise Exception(f"Error when revising PPTX: {e}")

        else:
            raise Exception(f"File type not supported : {file_type}")

        shutil.rmtree(temp_folder, ignore_errors=True)

        return response

    except Exception as e:
        shutil.rmtree(temp_folder, ignore_errors=True)
        return json.dumps(
            {"error": {"message": str(e)}},
            indent=4,
            ensure_ascii=False
        )

@mcp.tool()
async def create_file(data: dict, persistent: bool = PERSISTENT_FILES) -> dict:
    """
    Create a single file based on 'data' description.
    data examples:
      {"format":"pdf","filename":"report.pdf","content":[...]}
      {"format":"docx","filename":"doc.docx","content":[...],"title":"..."}
      {"format":"pptx","filename":"slides.pptx","slides_data":[...],"title":"..."}
      {"format":"xlsx","filename":"data.xlsx","content":[[...]],"title":"..."}
      {"format":"csv","filename":"data.csv","content":[[...]]}
      {"format":"txt|xml|py|...","filename":"file.ext","content":"string"}
    """
    log.debug("Creating file via tool (server.py)")
    folder_path = _generate_unique_folder()
    format_type = (data.get("format") or "").lower()
    filename = data.get("filename")
    content = data.get("content")
    title = data.get("title")

    if format_type == "pdf":
        result = create_pdf(content if isinstance(content, list) else [str(content or "")], filename, folder_path=folder_path)
    elif format_type == "pptx":
        result = create_presentation(
            data.get("slides_data", []),
            filename,
            folder_path=folder_path,
            title=title,
            pptx_template_path=PPTX_TEMPLATE_PATH,
        )
    elif format_type == "docx":
        result = create_word(
            content if content is not None else [],
            filename,
            folder_path=folder_path,
            title=title,
            docx_template_path=DOCX_TEMPLATE_PATH,
        )
    elif format_type == "xlsx":
        result = create_excel(
            content if content is not None else [],
            filename,
            folder_path=folder_path,
            title=title,
            xlsx_template_path=XLSX_TEMPLATE_PATH if "xlsx_template_path" in create_excel.__code__.co_varnames else None,  # type: ignore
        )
    elif format_type == "csv":
        result = _create_csv(content if content is not None else [], filename, folder_path=folder_path)
    else:
        use_filename = filename or f"export.{format_type or 'txt'}"
        result = _create_raw_file(content if content is not None else "", use_filename, folder_path=folder_path)

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": result["url"]}


@mcp.tool()
async def generate_and_archive(
    files_data: list[dict],
    archive_format: str = "zip",
    archive_name: str | None = None,
    persistent: bool = PERSISTENT_FILES
) -> dict:
    """
    Generate multiple files then archive them.
    files_data: list of 'data' dicts (same shape as for create_file)
    archive_format: zip | 7z | tar.gz
    """
    log.debug("Generating archive via tool (server.py)")
    folder_path = _generate_unique_folder()
    generated_paths: list[str] = []

    for file_info in files_data or []:
        fmt = (file_info.get("format") or "").lower()
        fname = file_info.get("filename")
        content = file_info.get("content")
        title = file_info.get("title")
        try:
            if fmt == "pdf":
                res = create_pdf(content if isinstance(content, list) else [str(content or "")], fname, folder_path=folder_path)
            elif fmt == "pptx":
                res = create_presentation(
                    file_info.get("slides_data", []),
                    fname,
                    folder_path=folder_path,
                    title=title,
                    pptx_template_path=PPTX_TEMPLATE_PATH,
                )
            elif fmt == "docx":
                res = create_word(
                    content if content is not None else [],
                    fname,
                    folder_path=folder_path,
                    title=title,
                    docx_template_path=DOCX_TEMPLATE_PATH,
                )
            elif fmt == "xlsx":
                res = create_excel(
                    content if content is not None else [],
                    fname,
                    folder_path=folder_path,
                    title=title,
                    xlsx_template_path=XLSX_TEMPLATE_PATH if "xlsx_template_path" in create_excel.__code__.co_varnames else None,  # type: ignore
                )
            elif fmt == "csv":
                res = _create_csv(content if content is not None else [], fname, folder_path=folder_path)
            else:
                use_fname = fname or f"export.{fmt or 'txt'}"
                res = _create_raw_file(content if content is not None else "", use_fname, folder_path=folder_path)
        except Exception as e:
            log.error(f"Error generating file {fname or '<no name>'}: {e}", exc_info=True)
            raise
        generated_paths.append(res["path"])

    timestamp = __import__("datetime").datetime.now().strftime("%Y%m%d_%H%M%S")
    archive_basename = f"{archive_name or 'archive'}_{timestamp}"
    if archive_format.lower() == "7z":
        archive_filename = f"{archive_basename}.7z"
    elif archive_format.lower() == "tar.gz":
        archive_filename = f"{archive_basename}.tar.gz"
    else:
        archive_filename = f"{archive_basename}.zip"

    archive_path = os.path.join(folder_path, archive_filename)

    if archive_format.lower() == "7z":
        with py7zr.SevenZipFile(archive_path, mode="w") as archive:
            for p in generated_paths:
                archive.write(p, os.path.relpath(p, folder_path))
    elif archive_format.lower() == "tar.gz":
        with tarfile.open(archive_path, "w:gz") as tar:
            for p in generated_paths:
                tar.add(p, arcname=os.path.relpath(p, folder_path))
    else:
        with zipfile.ZipFile(archive_path, "w") as zipf:
            for p in generated_paths:
                zipf.write(p, os.path.relpath(p, folder_path))

    if not persistent:
        _cleanup_files(folder_path, FILES_DELAY)

    return {"url": _public_url(folder_path, archive_filename)}

from sse_starlette.sse import EventSourceResponse

class SimpleRequestContext:
    def __init__(self, request):
        self.request = request

class SimpleCtx:
    def __init__(self, request):
        self.request_context = SimpleRequestContext(request)

async def handle_sse(request: Request) -> Response:
    """Handle SSE transport for MCP - supports both GET and POST"""
    
    if request.method == "POST":
        try:
            message = await request.json()
            log.debug(f"Received POST message: {message}")
            
            response = {
                "jsonrpc": "2.0",
                "id": message.get("id"),
                "result": None
            }
            
            method = message.get("method")
            
            if method == "initialize":
                response["result"] = {
                    "protocolVersion": "2024-11-05",
                    "capabilities": {
                        "tools": {},
                        "logging": {}
                    },
                    "serverInfo": {
                        "name": "file_export_mcp",
                        "version": SCRIPT_VERSION
                    }
                }
            elif method == "tools/list":
                response["result"] = {
                    "tools": [
                        {
                            "name": "create_file",
                            "description": "Create files in various formats (pdf, docx, pptx, xlsx, csv, txt, xml, py, etc.). Supports rich content including titles, paragraphs, lists, tables, images via queries, and more.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "data": {
                                        "type": "object",
                                        "description": "File data configuration",
                                        "properties": {
                                            "format": {
                                                "type": "string",
                                                "enum": ["pdf", "docx", "pptx", "xlsx", "csv", "txt", "xml", "py", "json", "md"],
                                                "description": "Output file format"
                                            },
                                            "filename": {
                                                "type": "string",
                                                "description": "Name of the file to create (optional, will be auto-generated if not provided)"
                                            },
                                            "title": {
                                                "type": "string",
                                                "description": "Document title (for docx, pptx, xlsx, pdf)"
                                            },
                                            "content": {
                                                "description": "Content varies by format. For pdf/docx: array (objects or strings). For xlsx/csv: 2D array. For pptx: use slides_data instead. For txt/xml/py: string",
                                                "oneOf": [
                                                    {
                                                        "type": "array",
                                                        "items": {
                                                            "anyOf": [
                                                                { "type": "string" },
                                                                { "type": "number" },
                                                                { "type": "boolean" },
                                                                { "type": "object" },
                                                                {
                                                                    "type": "array",
                                                                    "items": {
                                                                        "anyOf": [
                                                                            { "type": "string" },
                                                                            { "type": "number" },
                                                                            { "type": "boolean" },
                                                                            { "type": "object" },
                                                                            { "type": "null" }
                                                                        ]
                                                                    }
                                                                }
                                                            ]
                                                        }
                                                    },
                                                    { "type": "string" },
                                                    { "type": "object" },
                                                    { "type": "null" }
                                                ]
                                            },
                                            "slides_data": {
                                                "type": "array",
                                                "description": "For pptx format only: array of slide objects",
                                                "items": {
                                                    "type": "object",
                                                    "properties": {
                                                        "title": {"type": "string"},
                                                        "content": {
                                                            "type": "array",
                                                            "items": {"type": "string"}
                                                        },
                                                        "image_query": {
                                                            "type": "string",
                                                            "description": "Search query for image (Unsplash, Pexels, or local SD)"
                                                        },
                                                        "image_position": {
                                                            "type": "string",
                                                            "enum": ["left", "right", "top", "bottom"],
                                                            "description": "Position of the image on the slide"
                                                        },
                                                        "image_size": {
                                                            "type": "string",
                                                            "enum": ["small", "medium", "large"],
                                                            "description": "Size of the image"
                                                        }
                                                    }
                                                }
                                            }
                                        },
                                        "required": ["format"]
                                    },
                                    "persistent": {
                                        "type": "boolean",
                                        "description": "Whether to keep files permanently (default: false, files deleted after delay)"
                                    }
                                },
                                "required": ["data"]
                            }
                        },
                        {
                            "name": "generate_and_archive",
                            "description": "Generate multiple files and create an archive (zip, 7z, tar.gz)",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "files_data": {
                                        "type": "array",
                                        "description": "Array of file data objects",
                                        "items": {
                                            "type": "object",
                                            "properties": {
                                                "format": { "type": "string" },
                                                "filename": { "type": "string" },
                                                "content": {
                                                    "description": "For pdf/docx: array (objects or strings). For xlsx/csv: 2D array. For others: string/object",
                                                    "oneOf": [
                                                        {
                                                            "type": "array",
                                                            "items": {
                                                                "anyOf": [
                                                                    { "type": "string" },
                                                                    { "type": "number" },
                                                                    { "type": "boolean" },
                                                                    { "type": "object" },
                                                                    {
                                                                        "type": "array",
                                                                        "items": {
                                                                            "anyOf": [
                                                                                { "type": "string" },
                                                                                { "type": "number" },
                                                                                { "type": "boolean" },
                                                                                { "type": "object" },
                                                                                { "type": "null" }
                                                                            ]
                                                                        }
                                                                    }
                                                                ]
                                                            }
                                                        },
                                                        { "type": "string" },
                                                        { "type": "object" },
                                                        { "type": "null" }
                                                    ]
                                                },
                                                "title": { "type": "string" },
                                                "slides_data": {
                                                    "type": "array",
                                                    "description": "For pptx format only: array of slide objects",
                                                    "items": {
                                                        "type": "object",
                                                        "properties": {
                                                            "title": { "type": "string" },
                                                            "content": {
                                                                "type": "array",
                                                                "items": { "type": "string" }
                                                            },
                                                            "image_query": {
                                                                "type": "string",
                                                                "description": "Search query for image (Unsplash, Pexels, or local SD)"
                                                            },
                                                            "image_position": {
                                                                "type": "string",
                                                                "enum": ["left", "right", "top", "bottom"],
                                                                "description": "Position of the image on the slide"
                                                            },
                                                            "image_size": {
                                                                "type": "string",
                                                                "enum": ["small", "medium", "large"],
                                                                "description": "Size of the image"
                                                            }
                                                        }
                                                    }
                                                }
                                            },
                                            "required": ["format"]
                                        }
                                    },
                                    "archive_format": {"type": "string", "enum": ["zip", "7z", "tar.gz"]},
                                    "archive_name": {"type": "string"},
                                    "persistent": {"type": "boolean"}
                                },
                                "required": ["files_data"]
                            }
                        },
                        {
                            "name": "full_context_document",
                            "description": "Extract and return the complete structure, content, and metadata of a document (docx, xlsx, pptx). Returns a JSON structure with indexed elements (paragraphs, headings, tables, cells, slides, images) that can be referenced for editing or review.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI file upload"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension (e.g., 'report.docx', 'data.xlsx', 'presentation.pptx')"
                                    }
                                },
                                "required": ["file_id", "file_name"]
                            }
                        },
                        {
                            "name": "edit_document",
                            "description": "Edit an existing document (docx, xlsx, pptx) using structured operations. Supports inserting/deleting elements and updating content. ALWAYS call full_context_document() first to get proper IDs and references. Preserves formatting and returns a download link for the edited file.",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                    "edits": {
                                        "type": "object",
                                        "description": "Edit operations and content changes",
                                        "properties": {
                                            "ops": {
                                                "type": "array",
                                                "description": "Structural operations (insert/delete). For PPTX: ['insert_after', slide_id, 'nK'], ['insert_before', slide_id, 'nK'], ['delete_slide', slide_id]. For DOCX: ['insert_after', para_xml_id, 'nK'], ['insert_before', para_xml_id, 'nK'], ['delete_paragraph', para_xml_id]. For XLSX: ['insert_row', 'sheet_name', row_idx], ['delete_row', 'sheet_name', row_idx], ['insert_column', 'sheet_name', col_idx], ['delete_column', 'sheet_name', col_idx]",
                                                "items": {
                                                    "type": "array",
                                                    "items": {
                                                        "oneOf": [
                                                            {"type": "string"},
                                                            {"type": "integer"}
                                                        ]
                                                    }
                                                }
                                            },
                                            "content_edits": {
                                                "type": "array",
                                                "description": "Content updates. Prefer object items: {target, value}. For PPTX: target 'sid:<slide_id>/shid:<shape_id>' or 'nK:slot:title'/'body'/'table'. For DOCX: 'pid:<para_xml_id>' or 'tid:<table_xml_id>/cid:<cell_xml_id>' or 'nK'. For XLSX: 'A1', 'B5'.",
                                                "items": {
                                                    "type": "object",
                                                    "required": ["target", "value"],
                                                    "properties": {
                                                        "target": {
                                                            "type": "string",
                                                            "description": "Target reference (element ID or cell ref)"
                                                        },
                                                        "value": {
                                                            "description": "New content (string, number, boolean, array of strings, or 2D array for tables)",
                                                            "oneOf": [
                                                                {"type": "string"},
                                                                {"type": "number"},
                                                                {"type": "boolean"},
                                                                {"type": "array", "items": {"type": "string"}},
                                                                {
                                                                    "type": "array",
                                                                    "items": {
                                                                        "type": "array",
                                                                        "items": {
                                                                            "oneOf": [
                                                                                {"type": "string"},
                                                                                {"type": "number"},
                                                                                {"type": "boolean"},
                                                                                {"type": "null"}
                                                                            ]
                                                                        }
                                                                    }
                                                                }
                                                            ]
                                                        }
                                                    },
                                                    "additionalProperties": False
                                                }
                                            }
                                        }
                                    }
                                },
                                "required": ["file_id", "file_name", "edits"]
                            }
                        },
                        {
                            "name": "review_document",
                            "description": "Review and add comments/corrections to an existing document (docx, xlsx, pptx). Returns a download link for the reviewed document with comments added. For Excel, the index MUST be a cell reference (e.g., 'A1', 'B5', 'C10') as returned by full_context_document. For Word: use either an integer paragraph index or 'pid:<para_xml_id>'. For PowerPoint: use either an integer slide index or 'sid:<slide_id>' (optionally 'sid:<slide_id>/shid:<shape_id>' to target a shape).",
                            "inputSchema": {
                                "type": "object",
                                "properties": {
                                    "file_id": {
                                        "type": "string",
                                        "description": "The file ID from OpenWebUI"
                                    },
                                    "file_name": {
                                        "type": "string",
                                        "description": "The name of the file with extension"
                                    },
                                            "review_comments": {
                                                "type": "array",
                                                "description": "Array of objects {index, comment}. For Excel: index must be a cell reference string like 'A1', 'B3'. For Word: integer paragraph index or 'pid:<para_xml_id>'. For PowerPoint: integer slide index or 'sid:<slide_id>' (optionally 'sid:<slide_id>/shid:<shape_id>').",
                                                "items": {
                                                    "type": "object",
                                                    "required": ["index", "comment"],
                                                    "properties": {
                                                        "index": {
                                                            "description": "Index/reference: For Excel use cell reference (e.g., 'A1'); for Word/PowerPoint use integer or an id key string like 'pid:<para_xml_id>' / 'sid:<slide_id>'",
                                                            "oneOf": [
                                                                {"type": "string"},
                                                                {"type": "integer"}
                                                            ]
                                                        },
                                                        "comment": {
                                                            "type": "string",
                                                            "description": "Comment or correction text"
                                                        }
                                                    },
                                                    "additionalProperties": False
                                                }
                                            }
                                },
                                "required": ["file_id", "file_name", "review_comments"]
                            }
                        }
                    ]
                }
            elif method == "tools/call":
                params = message.get("params", {})
                tool_name = params.get("name")
                arguments = params.get("arguments", {}) or {}
                ctx = SimpleCtx(request)

                try:
                    if tool_name == "create_file":
                        result = await create_file(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": f"File created successfully: {result.get('url')}"
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "generate_and_archive":
                        result = await generate_and_archive(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": f"Archive created successfully: {result.get('url')}"
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "full_context_document":
                        arguments.setdefault("ctx", ctx)
                        result = await full_context_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": result
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "edit_document":
                        arguments.setdefault("ctx", ctx)
                        try:
                            edits_arg = arguments.get("edits")
                            if isinstance(edits_arg, dict):
                                ce = edits_arg.get("content_edits")
                                if isinstance(ce, list) and (len(ce) == 0 or isinstance(ce[0], dict)):
                                    edits_arg["content_edits"] = [
                                        [item.get("target"), item.get("value")]
                                        for item in (ce or [])
                                        if isinstance(item, dict) and "target" in item and "value" in item
                                    ]
                        except Exception:
                            pass
                        result = await edit_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps(result, indent=2, ensure_ascii=False)
                                }
                            ],
                            "isError": False
                        }

                    elif tool_name == "review_document":
                        arguments.setdefault("ctx", ctx)
                        try:
                            rc = arguments.get("review_comments")
                            if isinstance(rc, list) and (len(rc) == 0 or isinstance(rc[0], dict)):
                                arguments["review_comments"] = [
                                    [item.get("index"), item.get("comment")]
                                    for item in (rc or [])
                                    if isinstance(item, dict) and "index" in item and "comment" in item
                                ]
                        except Exception:
                            pass
                        result = await review_document(**arguments)
                        response["result"] = {
                            "content": [
                                {
                                    "type": "text",
                                    "text": json.dumps(result, indent=2, ensure_ascii=False)
                                }
                            ],
                            "isError": False
                        }

                    else:
                        response["error"] = {
                            "code": -32601,
                            "message": f"Tool not found: {tool_name}"
                        }
                except Exception as e:
                    log.error(f"Error executing tool {tool_name}: {e}", exc_info=True)
                    response["result"] = {
                        "content": [
                            {
                                "type": "text",
                                "text": f"Error executing tool: {str(e)}"
                            }
                        ],
                        "isError": True
                    }
            else:
                response["error"] = {
                    "code": -32601,
                    "message": f"Method not found: {method}"
                }
            
            return JSONResponse(response)
            
        except Exception as e:
            log.error(f"Error handling POST request: {e}", exc_info=True)
            return JSONResponse(
                {
                    "jsonrpc": "2.0",
                    "id": None,
                    "error": {
                        "code": -32700,
                        "message": f"Parse error: {str(e)}"
                    }
                },
                status_code=400
            )
    
    else:
        async def event_generator():
            """Generator for SSE events with correct text format"""
            try:
                endpoint_data = json.dumps({"endpoint": "/sse"})
                yield f"event: endpoint\ndata: {endpoint_data}\n\n"
                
                import asyncio
                while True:
                    await asyncio.sleep(15)
                    yield f"event: ping\ndata: {{}}\n\n"
                    
            except asyncio.CancelledError:
                log.info("SSE connection closed by client")
                raise
            except Exception as e:
                log.error(f"SSE Error: {e}", exc_info=True)
                error_data = json.dumps({"error": str(e)})
                yield f"event: error\ndata: {error_data}\n\n"
        
        return EventSourceResponse(
            event_generator(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "X-Accel-Buffering": "no",
                "Connection": "keep-alive"
            }
        )

async def handle_messages(request: Request) -> Response:
    """Handle POST requests to /messages endpoint"""
    try:
        data = await request.json()
        log.debug(f"Received message: {data}")
        
        response = {
            "jsonrpc": "2.0",
            "id": data.get("id"),
            "result": {
                "content": [
                    {
                        "type": "text",
                        "text": "Message received"
                    }
                ]
            }
        }
        
        return JSONResponse(response)
    except Exception as e:
        log.error(f"Message handling error: {e}", exc_info=True)
        return JSONResponse(
            {
                "jsonrpc": "2.0",
                "id": None,
                "error": {
                    "code": -32700,
                    "message": f"Parse error: {str(e)}"
                }
            },
            status_code=400
        )

async def health_check(request: Request) -> Response:
    """Health check endpoint"""
    return JSONResponse({"status": "healthy", "server": "file_export_mcp"})

app = Starlette(
    debug=True,
    routes=[
        Route("/sse", endpoint=handle_sse, methods=["GET", "POST"]),
        Route("/messages", endpoint=handle_messages, methods=["POST"]),
        Route("/health", endpoint=health_check, methods=["GET"]),
    ]
)

if __name__ == "__main__":

    mode = (os.getenv("MODE", "SSE"))
 
    if mode == "sse":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
            
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in SSE mode on http://{host}:{port}")
        log.info(f"SSE endpoint: http://{host}:{port}/sse")
        log.info(f"Messages endpoint: http://{host}:{port}/messages")
            
        uvicorn.run(
            app,
            host=host,
            port=port,
            access_log=False,
            log_level="info",
            use_colors=False
        )
    elif mode == "http":
        port = int(os.getenv("MCP_HTTP_PORT", "9004"))
        host = os.getenv("MCP_HTTP_HOST", "0.0.0.0")
        
        log.info(f"Starting file_export_mcp version {SCRIPT_VERSION}")
        log.info(f"Starting file_export_mcp in http mode on http://{host}:{port}")
        log.info(f"HTTP endpoint: http://{host}:{port}/mcp")

        mcp.run(
            transport="streamable-http"
        )